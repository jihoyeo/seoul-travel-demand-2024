"""02_sgis_census.pptx — M2 슬라이드 (16:9, MiniMax 디자인 시스템)."""
from __future__ import annotations

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from slide_theme import (
    C, SLIDE_W, SLIDE_H, MARGIN_X, MARGIN_Y, CONTENT_W, CONTENT_H,
    add_text, add_rounded, add_card_white, add_card_surface, add_card_vibrant,
    add_pill_button, add_badge, add_hairline, add_page_chrome,
    add_code_block,
    new_blank_slide, setup_169, radius_for, TYPE,
)

MODULE_LABEL = "M2  ·  SGIS 인구·가구·주택 통계"
DECK_LABEL = "MiniMax Studio  ·  TAZ Lecture  ·  M2"
ACCENT = C.magenta   # M2 brand color
TOTAL = 11


# ---------- slide 1: hero title ----------
def slide_title(prs, idx):
    s = new_blank_slide(prs)

    # right vibrant card
    coral_w = Inches(5.0)
    add_rounded(s, SLIDE_W - coral_w - Inches(0.4), Inches(0.4),
                coral_w, SLIDE_H - Inches(0.8),
                fill=ACCENT, line=None,
                radius_ratio=radius_for(coral_w, SLIDE_H - Inches(0.8), 0.42))

    # promo strip
    promo_w = Inches(4.5)
    promo = add_rounded(s, MARGIN_X, Inches(0.45), promo_w, Inches(0.42),
                        fill=C.primary, line=None, radius_ratio=0.5)
    tf = promo.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    from slide_theme import _set_para, _set_font
    _set_para(p, align=PP_ALIGN.CENTER, line_spacing=1.0)
    r = p.add_run(); r.text = "LECTURE  ·  M2 of 7  ·  SGIS Census"
    _set_font(r, size_pt=10, color=C.on_primary, bold=True, letter_spacing_pt=0.8)

    # hero
    add_text(s, MARGIN_X, Inches(1.7), Inches(7.2), Inches(2.6),
             [("인구 · 가구 ·\n주택 통계",
               {"style": "hero", "color": C.ink, "letter_spacing_pt": -1.8,
                "line_spacing": 1.05})])

    # subtitle
    add_text(s, MARGIN_X, Inches(4.6), Inches(7.2), Inches(1.0),
             [("OA 단위 통계 10종 → wide pivot → 분석 가능 테이블.",
               {"style": "subtitle", "color": C.slate})])

    # buttons
    btn_y = Inches(5.55)
    add_pill_button(s, MARGIN_X, btn_y, Inches(2.0), Inches(0.5),
                    "Open Notebook  →", primary=True)
    add_pill_button(s, MARGIN_X + Inches(2.15), btn_y, Inches(1.7), Inches(0.5),
                    "View Repo", primary=False)

    # right card content (on magenta)
    pad = Inches(0.45)
    add_text(s, SLIDE_W - coral_w + pad, Inches(0.95),
             coral_w - pad * 2, Inches(2.4),
             [("M2", {"style": "hero", "color": C.on_dark,
                       "letter_spacing_pt": -2.5, "line_spacing": 1.0,
                       "size_pt": 96})])
    add_text(s, SLIDE_W - coral_w + pad, Inches(3.5),
             coral_w - pad * 2, Inches(0.5),
             [("SGIS CENSUS", {"style": "caption_b",
                                "color": C.on_dark, "letter_spacing_pt": 1.5})])
    bullets = [
        ("• 11 CSV (cp949, no header)", {"style": "body_md", "color": C.on_dark}),
        ("• 19,097 OA × ~140 컬럼", {"style": "body_md", "color": C.on_dark, "space_before": 4}),
        ("• 5인 미만 N/A 보호 셀", {"style": "body_md", "color": C.on_dark, "space_before": 4}),
        ("• 426 행정동 집계 가능", {"style": "body_md", "color": C.on_dark, "space_before": 4}),
    ]
    add_text(s, SLIDE_W - coral_w + pad, Inches(4.2),
             coral_w - pad * 2, Inches(2.6), bullets)

    add_badge(s, SLIDE_W - coral_w + pad, Inches(0.6), "M2 · 2024", kind="new")

    # bottom chrome
    add_text(s, MARGIN_X, SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             [(DECK_LABEL, {"style": "micro", "color": C.steel})])
    add_text(s, SLIDE_W - MARGIN_X - Inches(2), SLIDE_H - Inches(0.45),
             Inches(2), Inches(0.3),
             [(f"{idx:02d} / {TOTAL:02d}",
               {"style": "micro", "color": C.steel, "letter_spacing_pt": 0.5})],
             align=PP_ALIGN.RIGHT)


# ---------- slide 2: 학습 목표 ----------
def slide_objectives(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(1.2),
             [("학습 목표",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.0), Inches(11), Inches(0.5),
             [("Five things you'll be able to do by the end of this module.",
               {"style": "subtitle", "color": C.slate})])

    objectives = [
        ("01", "long → wide pivot 익히기",
         "SGIS 통계 4-컬럼 long 스키마와 pivot_table 패턴."),
        ("02", "변수 코드 8 그룹 정리",
         "to_in_* / in_age_* / to_ga_* / ga_sd_* / to_ho_* / ho_gb_* / ho_ar_* / ho_yr_*"),
        ("03", "N/A 통계 보호 보존",
         "5인 미만 셀은 NaN 으로 — 0 imputation 금지의 이유."),
        ("04", "OA 단위 choropleth 4종",
         "인구밀도 · 노령화 · 아파트비율 · 평균가구원수."),
        ("05", "행정동 단위 집계",
         "TOT_OA_CD[:8] groupby 로 426 행정동 합계."),
    ]

    card_w = (CONTENT_W - Inches(0.4) * 2) / 3
    card_h = Inches(1.85)
    gap = Inches(0.25)
    start_y = Inches(2.85)

    for i, (no, title, desc) in enumerate(objectives):
        col = i % 3
        row = i // 3
        x = MARGIN_X + (card_w + gap) * col
        y = start_y + (card_h + gap) * row
        add_card_white(s, x, y, card_w, card_h, radius_in=0.21)
        add_text(s, x + Inches(0.25), y + Inches(0.2),
                 Inches(0.7), Inches(0.3),
                 [(no, {"style": "caption_b", "color": ACCENT, "letter_spacing_pt": 0.8})])
        add_text(s, x + Inches(0.25), y + Inches(0.5),
                 card_w - Inches(0.5), Inches(0.55),
                 [(title, {"style": "card_title", "color": C.ink})])
        add_text(s, x + Inches(0.25), y + Inches(1.05),
                 card_w - Inches(0.5), Inches(0.7),
                 [(desc, {"style": "body_sm", "color": C.slate})])


# ---------- slide 3: 데이터 메타 ----------
def slide_data_meta(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("DATA · SGIS Census",
               {"style": "caption_b", "color": ACCENT, "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("data/raw/sgis_census/",
               {"style": "heading_lg", "color": C.ink, "code": True,
                "letter_spacing_pt": -0.5})])

    rows = [
        ("출처",         "통계청 SGIS — 집계구 단위 통계"),
        ("기준",         "2024년 (인구주택총조사 등록센서스)"),
        ("형식",         "CSV 10개 · 인코딩 cp949 · 헤더 없음"),
        ("키",           "TOT_OA_CD (14자리) — 2025 SGIS OA 와 동일"),
        ("통합 스키마",  "year, TOT_OA_CD, var_code, value (str)"),
    ]
    table_x = MARGIN_X
    table_y = Inches(2.55)
    table_w = Inches(7.5)
    row_h = Inches(0.55)

    add_card_white(s, table_x, table_y,
                   table_w, row_h * len(rows) + Inches(0.3),
                   radius_in=0.21)
    for i, (k, v) in enumerate(rows):
        ry = table_y + Inches(0.15) + row_h * i
        add_text(s, table_x + Inches(0.5), ry, Inches(2.0), row_h,
                 [(k, {"style": "body_sm_med", "color": C.steel})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, table_x + Inches(2.6), ry, table_w - Inches(2.9), row_h,
                 [(v, {"style": "body_md", "color": C.ink})],
                 anchor=MSO_ANCHOR.MIDDLE)
        if i < len(rows) - 1:
            add_hairline(s, table_x + Inches(0.5),
                         ry + row_h - Inches(0.01),
                         table_w - Inches(1.0),
                         color=C.hairline_soft, weight_pt=0.5)

    # right: schema callout (vibrant)
    cx = table_x + table_w + Inches(0.3)
    cw = SLIDE_W - MARGIN_X - cx
    cy = table_y
    ch = row_h * len(rows) + Inches(0.3)
    add_card_vibrant(s, cx, cy, cw, ch, color=ACCENT, radius_in=0.42)
    add_text(s, cx + Inches(0.4), cy + Inches(0.4),
             cw - Inches(0.8), Inches(0.35),
             [("KEY", {"style": "caption_b", "color": C.on_dark,
                       "letter_spacing_pt": 1.2})])
    add_text(s, cx + Inches(0.4), cy + Inches(0.85),
             cw - Inches(0.8), Inches(1.6),
             [("TOT_OA_CD\n14 자리",
               {"style": "heading_md", "color": C.on_dark, "size_pt": 28,
                "code": True, "line_spacing": 1.15, "letter_spacing_pt": -0.5})])
    add_text(s, cx + Inches(0.4), cy + Inches(2.25),
             cw - Inches(0.8), Inches(0.7),
             [("이 키로 모든 통계가\nM4 oa_master 에 합쳐진다.",
               {"style": "body_sm", "color": C.on_dark, "line_spacing": 1.55})])


# ---------- slide 4: N/A 보호 (강조 슬라이드) ----------
def slide_na_protection(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    # full-bleed magenta card
    cy = Inches(1.05)
    ch = Inches(5.6)
    add_card_vibrant(s, MARGIN_X, cy, CONTENT_W, ch, color=ACCENT, radius_in=0.42)

    add_text(s, MARGIN_X + Inches(0.6), cy + Inches(0.5),
             CONTENT_W - Inches(1.2), Inches(0.4),
             [("⚠  PITFALL · 통계 보호",
               {"style": "caption_b", "color": C.on_dark,
                "letter_spacing_pt": 1.5})])

    # 좌: 거대 헤드
    add_text(s, MARGIN_X + Inches(0.6), cy + Inches(1.05),
             Inches(7.3), Inches(2.6),
             [("왜 value 가\nstring 인가",
               {"style": "hero", "color": C.on_dark, "size_pt": 60,
                "letter_spacing_pt": -1.8, "line_spacing": 1.10})])
    add_text(s, MARGIN_X + Inches(0.6), cy + Inches(3.7),
             Inches(7.3), Inches(1.5),
             [("통계청은 5인 미만 셀을 N/A 로 마스킹한다.\n"
               "pd.to_numeric(errors='coerce') 로 그대로 NaN 보존.",
               {"style": "subtitle", "color": C.on_dark, "line_spacing": 1.55,
                "size_pt": 16})])

    # 우: 금지 룰 카드
    rx = MARGIN_X + Inches(7.95)
    rw = CONTENT_W - Inches(7.95) - Inches(0.6)
    ry_top = cy + Inches(1.05)

    # warning card
    add_rounded(s, rx, ry_top, rw, Inches(2.0),
                fill=C.canvas, line=None,
                radius_ratio=radius_for(rw, Inches(2.0), 0.32))
    add_badge(s, rx + Inches(0.3), ry_top + Inches(0.3), "DON'T", kind="new")
    add_text(s, rx + Inches(0.3), ry_top + Inches(0.85),
             rw - Inches(0.6), Inches(0.6),
             [("0 imputation\n금지",
               {"style": "heading_md", "color": C.ink, "size_pt": 22,
                "line_spacing": 1.15, "letter_spacing_pt": -0.3})])
    add_text(s, rx + Inches(0.3), ry_top + Inches(1.5),
             rw - Inches(0.6), Inches(0.4),
             [("0 으로 채우면 인구 0 인 OA 와\n보호된 OA 가 구별 안 된다.",
               {"style": "body_sm", "color": C.slate, "line_spacing": 1.45,
                "size_pt": 9})])

    # do card
    ry_bot = ry_top + Inches(2.15)
    add_rounded(s, rx, ry_bot, rw, Inches(2.0),
                fill=C.canvas, line=None,
                radius_ratio=radius_for(rw, Inches(2.0), 0.32))
    add_badge(s, rx + Inches(0.3), ry_bot + Inches(0.3), "DO", kind="success")
    add_text(s, rx + Inches(0.3), ry_bot + Inches(0.85),
             rw - Inches(0.6), Inches(0.6),
             [("NaN 보존",
               {"style": "heading_md", "color": C.ink, "size_pt": 22,
                "letter_spacing_pt": -0.3})])
    add_text(s, rx + Inches(0.3), ry_bot + Inches(1.4),
             rw - Inches(0.6), Inches(0.5),
             [("errors='coerce' 한 줄로\n끝. 분석은 NaN-aware 로.",
               {"style": "body_sm", "color": C.slate, "line_spacing": 1.45,
                "size_pt": 9})])


# ---------- slide 5: 등록센서스 vs 표본조사 ----------
def slide_register_vs_sample(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(11), Inches(1.2),
             [("등록센서스 vs 표본조사",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.0), Inches(11), Inches(0.5),
             [("본 강좌는 등록센서스만 사용 — 행정자료 기반 stock 데이터.",
               {"style": "subtitle", "color": C.slate})])

    cw = (CONTENT_W - Inches(0.4)) / 2
    ch = Inches(3.85)
    cy = Inches(2.85)

    # LEFT — 등록센서스 (vibrant blue, used)
    lx = MARGIN_X
    add_card_vibrant(s, lx, cy, cw, ch, color=C.blue_deep, radius_in=0.42)
    add_text(s, lx + Inches(0.5), cy + Inches(0.45),
             cw - Inches(1.0), Inches(0.35),
             [("USED · 본 강좌",
               {"style": "caption_b", "color": C.on_dark,
                "letter_spacing_pt": 1.5})])
    add_text(s, lx + Inches(0.5), cy + Inches(0.85),
             cw - Inches(1.0), Inches(0.7),
             [("등록센서스",
               {"style": "heading_lg", "color": C.on_dark,
                "letter_spacing_pt": -0.5})])
    add_text(s, lx + Inches(0.5), cy + Inches(1.55),
             cw - Inches(1.0), Inches(0.4),
             [("2015 도입 · 행정자료 기반 (주민등록 · 세대 · 건축물대장)",
               {"style": "body_sm", "color": C.on_dark})])
    bullets = [
        ("✓ 모든 가구 커버 (전수)",
         {"style": "body_md", "color": C.on_dark, "space_before": 6}),
        ("✓ stock · 시점 정확",
         {"style": "body_md", "color": C.on_dark, "space_before": 4}),
        ("✓ 5세 단위 연령 → 예측 모델 강력 입력",
         {"style": "body_md", "color": C.on_dark, "space_before": 4}),
        ("✗ 통근 통행 같은 행동은 측정 안 됨",
         {"style": "body_md", "color": C.on_dark, "space_before": 4}),
    ]
    add_text(s, lx + Inches(0.5), cy + Inches(2.05),
             cw - Inches(1.0), Inches(1.7), bullets)

    # RIGHT — 표본조사 (white muted, not used)
    rx = lx + cw + Inches(0.4)
    add_card_white(s, rx, cy, cw, ch, radius_in=0.21)
    add_text(s, rx + Inches(0.5), cy + Inches(0.45),
             cw - Inches(1.0), Inches(0.35),
             [("NOT USED",
               {"style": "caption_b", "color": C.steel,
                "letter_spacing_pt": 1.5})])
    add_text(s, rx + Inches(0.5), cy + Inches(0.85),
             cw - Inches(1.0), Inches(0.7),
             [("표본조사",
               {"style": "heading_lg", "color": C.ink,
                "letter_spacing_pt": -0.5})])
    add_text(s, rx + Inches(0.5), cy + Inches(1.55),
             cw - Inches(1.0), Inches(0.4),
             [("전수조사 일부 · 가중치 있는 추정치",
               {"style": "body_sm", "color": C.slate})])
    bullets2 = [
        ("· 표본 가중 추정",
         {"style": "body_md", "color": C.charcoal, "space_before": 6}),
        ("· 분산이 큰 카테고리 변수에 적합",
         {"style": "body_md", "color": C.charcoal, "space_before": 4}),
        ("· 본 프로젝트에서는 사용 X",
         {"style": "body_md", "color": C.charcoal, "space_before": 4}),
    ]
    add_text(s, rx + Inches(0.5), cy + Inches(2.05),
             cw - Inches(1.0), Inches(1.7), bullets2)


# ---------- slide 6: 변수 코드 그룹 ----------
def slide_var_groups(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("SCHEMA · 변수 코드 그룹 8종",
               {"style": "caption_b", "color": ACCENT, "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("ALIAS 매핑으로 인간이 읽기 쉽게",
               {"style": "display_lg", "color": C.ink, "size_pt": 36})])
    add_text(s, MARGIN_X, Inches(2.25), Inches(11), Inches(0.5),
             [("build_oa_master.py 의 ALIAS — to_in_001 → pop_total 식.",
               {"style": "subtitle", "color": C.slate})])

    groups = [
        ("인구총괄",  "to_in_*",   "총인구 · 남 · 여 · 인구밀도 · 노령화지수"),
        ("성연령",    "in_age_*",  "성×5세 (남자 0-9 = 10년, 외 5년) · 81 변수"),
        ("가구총괄",  "to_ga_*",   "가구수 · 평균가구원수"),
        ("세대구성",  "ga_sd_*",   "1·2·3·4세대 · 1인 · 비친족"),
        ("주택총괄",  "to_ho_*",   "총주택수"),
        ("주택유형",  "ho_gb_*",   "단독 · 아파트 · 연립 · 다세대 · 비거주용 · 기타"),
        ("연건평",    "ho_ar_*",   "9 면적구간"),
        ("건축년도",  "ho_yr_*",   "20 년대구간"),
    ]

    # 4 × 2 grid
    cw = (CONTENT_W - Inches(0.25) * 3) / 4
    ch = Inches(1.7)
    gy = Inches(2.95)
    for i, (name, code, desc) in enumerate(groups):
        col = i % 4
        row = i // 4
        x = MARGIN_X + (cw + Inches(0.25)) * col
        y = gy + (ch + Inches(0.22)) * row
        add_card_white(s, x, y, cw, ch, radius_in=0.21)
        # group name
        add_text(s, x + Inches(0.25), y + Inches(0.22),
                 cw - Inches(0.5), Inches(0.4),
                 [(name, {"style": "card_title", "color": C.ink, "size_pt": 14})])
        # code chip
        add_text(s, x + Inches(0.25), y + Inches(0.65),
                 cw - Inches(0.5), Inches(0.35),
                 [(code, {"style": "code", "color": ACCENT, "code": True,
                          "size_pt": 11, "bold": True})])
        # desc
        add_text(s, x + Inches(0.25), y + Inches(1.05),
                 cw - Inches(0.5), Inches(0.6),
                 [(desc, {"style": "body_sm", "color": C.slate, "size_pt": 9.5,
                          "line_spacing": 1.4})])


# ---------- slide 7: long → wide pivot ----------
def slide_pivot(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(11), Inches(1.2),
             [("long → wide pivot",
               {"style": "display_lg", "color": C.ink, "code": True,
                "letter_spacing_pt": -0.5})])
    add_text(s, MARGIN_X, Inches(2.0), Inches(11), Inches(0.5),
             [("11개 파일 ≈ 1.5 M long-row → 19,097 행 × ~140 열 wide.",
               {"style": "subtitle", "color": C.slate})])

    # left: code panel (dark)
    cw = Inches(8.4)
    ch = Inches(4.05)
    cy = Inches(2.75)
    py_lines = [
        ("parts = []",                                                     None),
        ("for fname in CENSUS_FILES:",                                     None),
        ("    df = pd.read_csv(p, header=None,",                            None),
        ("                     names=['year','TOT_OA_CD','var','value'],", None),
        ("                     encoding='cp949', dtype=str)",              None),
        ("    df['value'] = pd.to_numeric(df['value'], errors='coerce')",  "# 'N/A' → NaN"),
        ("    parts.append(df[['TOT_OA_CD','var','value']])",              None),
        ("",                                                                None),
        ("long = pd.concat(parts, ignore_index=True)",                     None),
        ("wide = long.pivot_table(index='TOT_OA_CD', columns='var',",      None),
        ("                        values='value', aggfunc='first')",       None),
    ]
    add_code_block(s, MARGIN_X, cy, cw, ch, py_lines,
                   dark=True, label="PYTHON · build_oa_master", show_dots=True,
                   font_pt=10, line_spacing=1.30, padding_in=0.36)

    # right: scale stats (vibrant magenta)
    rx = MARGIN_X + cw + Inches(0.3)
    rw = SLIDE_W - MARGIN_X - rx
    add_card_vibrant(s, rx, cy, rw, ch, color=ACCENT, radius_in=0.42)
    add_text(s, rx + Inches(0.4), cy + Inches(0.45),
             rw - Inches(0.8), Inches(0.35),
             [("SCALE", {"style": "caption_b", "color": C.on_dark,
                          "letter_spacing_pt": 1.2})])
    add_text(s, rx + Inches(0.4), cy + Inches(0.85),
             rw - Inches(0.8), Inches(0.95),
             [("≈ 1.5 M",
               {"style": "hero", "color": C.on_dark, "size_pt": 44,
                "letter_spacing_pt": -1.0, "line_spacing": 1.0})])
    add_text(s, rx + Inches(0.4), cy + Inches(1.85),
             rw - Inches(0.8), Inches(0.4),
             [("long-row · 11 files",
               {"style": "body_sm", "color": C.on_dark})])

    add_hairline(s, rx + Inches(0.4), cy + Inches(2.4),
                 rw - Inches(0.8), color=C.on_dark, weight_pt=0.5)

    add_text(s, rx + Inches(0.4), cy + Inches(2.55),
             rw - Inches(0.8), Inches(0.95),
             [("19,097 × 140",
               {"style": "hero", "color": C.on_dark, "size_pt": 30,
                "letter_spacing_pt": -0.5, "line_spacing": 1.0})])
    add_text(s, rx + Inches(0.4), cy + Inches(3.4),
             rw - Inches(0.8), Inches(0.4),
             [("rows × cols · pivot 후",
               {"style": "body_sm", "color": C.on_dark})])


# ---------- slide 8: 시각화 4종 ----------
def slide_visualizations(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("VIZ · choropleth 4종",
               {"style": "caption_b", "color": ACCENT, "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("OA 위에 자치구 경계 overlay",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.4), Inches(11), Inches(0.5),
             [("EPSG:5179 (m 단위) 에서 그리고 자치구 dissolve 로 overlay.",
               {"style": "subtitle", "color": C.slate})])

    vizs = [
        ("01", "인구밀도",          "log scale",  "강남 · 신촌 · 관악 hotspot",  C.coral),
        ("02", "노령화지수",         "비율",       "외곽일수록 ↑",                C.purple),
        ("03", "주택유형 — 아파트",   "% 비율",     "아파트 단지 분포 시각화",      C.blue),
        ("04", "가구 평균규모",      "평균",       "1인 가구 ↓ 영역 식별",          ACCENT),
    ]

    cw = (CONTENT_W - Inches(0.3) * 3) / 4
    ch = Inches(3.05)
    cy = Inches(3.0)
    for i, (no, title, scale, desc, color) in enumerate(vizs):
        x = MARGIN_X + (cw + Inches(0.3)) * i
        add_card_vibrant(s, x, cy, cw, ch, color=color, radius_in=0.42)
        add_text(s, x + Inches(0.35), cy + Inches(0.35),
                 cw - Inches(0.7), Inches(0.4),
                 [("VIZ", {"style": "caption_b", "color": C.on_dark,
                            "letter_spacing_pt": 1.2})])
        add_text(s, x + Inches(0.35), cy + Inches(0.7),
                 cw - Inches(0.7), Inches(1.0),
                 [(no, {"style": "hero", "color": C.on_dark,
                        "size_pt": 56, "letter_spacing_pt": -1.5,
                        "line_spacing": 1.0})])
        add_text(s, x + Inches(0.35), cy + Inches(1.65),
                 cw - Inches(0.7), Inches(0.4),
                 [(title, {"style": "card_title", "color": C.on_dark,
                            "size_pt": 15})])
        add_text(s, x + Inches(0.35), cy + Inches(2.05),
                 cw - Inches(0.7), Inches(0.35),
                 [(scale, {"style": "caption_b", "color": C.on_dark,
                            "letter_spacing_pt": 0.8})])
        add_text(s, x + Inches(0.35), cy + Inches(2.4),
                 cw - Inches(0.7), Inches(0.5),
                 [(desc, {"style": "body_sm", "color": C.on_dark,
                          "size_pt": 9.5, "line_spacing": 1.4})])


# ---------- slide 9: N/A 분석 ----------
def slide_na_analysis(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(11), Inches(1.2),
             [("어떤 변수에 보호 셀이 많은가",
               {"style": "display_lg", "color": C.ink, "size_pt": 36})])
    add_text(s, MARGIN_X, Inches(1.85), Inches(11), Inches(0.5),
             [("wide.isna().sum() 으로 변수별 결측 개수 확인.",
               {"style": "subtitle", "color": C.slate})])

    # 두 stat 카드 + code snippet
    cw = (CONTENT_W - Inches(0.4)) / 2
    ch = Inches(2.6)
    cy = Inches(2.75)

    # 좌: 거의 0 결측 (success green)
    lx = MARGIN_X
    add_card_vibrant(s, lx, cy, cw, ch, color=C.success_text, radius_in=0.42)
    add_text(s, lx + Inches(0.45), cy + Inches(0.4),
             cw - Inches(0.9), Inches(0.4),
             [("≈ 0 결측", {"style": "caption_b", "color": C.on_dark,
                              "letter_spacing_pt": 1.2})])
    add_text(s, lx + Inches(0.45), cy + Inches(0.85),
             cw - Inches(0.9), Inches(0.85),
             [("전수성 변수",
               {"style": "heading_md", "color": C.on_dark, "size_pt": 26,
                "letter_spacing_pt": -0.3})])
    add_text(s, lx + Inches(0.45), cy + Inches(1.7),
             cw - Inches(0.9), Inches(0.45),
             [("to_in_001 (총인구) · pop_density",
               {"style": "code", "color": C.on_dark, "code": True})])
    add_text(s, lx + Inches(0.45), cy + Inches(2.05),
             cw - Inches(0.9), Inches(0.4),
             [("개인정보 위험 없는 합계 변수.",
               {"style": "body_sm", "color": C.on_dark})])

    # 우: 결측 다수 (warn coral)
    rx = lx + cw + Inches(0.4)
    add_card_vibrant(s, rx, cy, cw, ch, color=C.coral, radius_in=0.42)
    add_text(s, rx + Inches(0.45), cy + Inches(0.4),
             cw - Inches(0.9), Inches(0.4),
             [("결측 다수", {"style": "caption_b", "color": C.on_dark,
                                "letter_spacing_pt": 1.2})])
    add_text(s, rx + Inches(0.45), cy + Inches(0.85),
             cw - Inches(0.9), Inches(0.85),
             [("세부 카테고리",
               {"style": "heading_md", "color": C.on_dark, "size_pt": 26,
                "letter_spacing_pt": -0.3})])
    add_text(s, rx + Inches(0.45), cy + Inches(1.7),
             cw - Inches(0.9), Inches(0.45),
             [("ga_sd_* · ho_gb_*",
               {"style": "code", "color": C.on_dark, "code": True})])
    add_text(s, rx + Inches(0.45), cy + Inches(2.05),
             cw - Inches(0.9), Inches(0.4),
             [("5인 미만 셀이 많아 결측 비율 ↑.",
               {"style": "body_sm", "color": C.on_dark})])

    # 하단 안내 (next pointer)
    py = cy + ch + Inches(0.3)
    ph = SLIDE_H - py - Inches(0.7)
    add_card_surface(s, MARGIN_X, py, CONTENT_W, ph, radius_in=0.16)
    add_text(s, MARGIN_X + Inches(0.4), py + Inches(0.2),
             CONTENT_W - Inches(0.8), Inches(0.35),
             [("NEXT · M6", {"style": "caption_b", "color": C.steel,
                              "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X + Inches(0.4), py + Inches(0.55),
             CONTENT_W - Inches(0.8), Inches(0.6),
             [("회귀 모델에서 NaN 처리 전략 — median imputation? 행 제거? mask 컬럼? — 을 결정해야 한다.",
               {"style": "body_md", "color": C.charcoal})])


# ---------- slide 10: 행정동 단위 집계 ----------
def slide_admdong_groupby(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(11), Inches(1.2),
             [("행정동 단위 집계",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.0), Inches(11), Inches(0.5),
             [("TOT_OA_CD 앞 8자리가 ADM_CD — groupby 한 줄.",
               {"style": "subtitle", "color": C.slate})])

    # left: code panel
    cw = Inches(8.4)
    ch = Inches(3.6)
    cy = Inches(2.85)
    py_lines = [
        ("wide['adm_cd'] = wide['TOT_OA_CD'].str[:8]",  None),
        ("adm = wide.groupby('adm_cd').sum()",          None),
    ]
    add_code_block(s, MARGIN_X, cy, cw, ch, py_lines,
                   dark=True, label="PYTHON", show_dots=True,
                   font_pt=14, line_spacing=1.45, padding_in=0.4)

    # right: result stat
    rx = MARGIN_X + cw + Inches(0.3)
    rw = SLIDE_W - MARGIN_X - rx
    add_card_vibrant(s, rx, cy, rw, ch, color=C.blue, radius_in=0.42)
    add_text(s, rx + Inches(0.4), cy + Inches(0.5),
             rw - Inches(0.8), Inches(0.4),
             [("RESULT", {"style": "caption_b", "color": C.on_dark,
                           "letter_spacing_pt": 1.2})])
    add_text(s, rx + Inches(0.4), cy + Inches(1.0),
             rw - Inches(0.8), Inches(1.5),
             [("426",
               {"style": "hero", "color": C.on_dark, "size_pt": 80,
                "letter_spacing_pt": -2.0, "line_spacing": 1.0})])
    add_text(s, rx + Inches(0.4), cy + Inches(2.4),
             rw - Inches(0.8), Inches(0.4),
             [("행정동",
               {"style": "card_title", "color": C.on_dark, "size_pt": 16})])
    add_text(s, rx + Inches(0.4), cy + Inches(2.85),
             rw - Inches(0.8), Inches(0.5),
             [("인구 · 가구 · 주택\n합계 자동 산출",
               {"style": "body_sm", "color": C.on_dark, "line_spacing": 1.45})])


# ---------- slide 11: 핵심 메시지 ----------
def slide_key_message(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    cy = Inches(1.05)
    ch = Inches(5.6)
    add_card_vibrant(s, MARGIN_X, cy, CONTENT_W, ch, color=ACCENT, radius_in=0.42)

    add_text(s, MARGIN_X + Inches(0.6), cy + Inches(0.55),
             CONTENT_W - Inches(1.2), Inches(0.4),
             [("KEY MESSAGE", {"style": "caption_b", "color": C.on_dark,
                                "letter_spacing_pt": 1.5})])

    add_text(s, MARGIN_X + Inches(0.6), cy + Inches(1.05),
             CONTENT_W - Inches(1.2), Inches(3.0),
             [("OA 키 한 곳에 모이는\n정적 stock 데이터 — 잘 정리해 두면\n거의 모든 후속 분석의 기본 입력이 된다.",
               {"style": "hero", "color": C.on_dark, "size_pt": 38,
                "letter_spacing_pt": -1.0, "line_spacing": 1.18})])

    add_text(s, MARGIN_X + Inches(0.6), cy + ch - Inches(1.4),
             CONTENT_W - Inches(1.2), Inches(0.5),
             [("N/A 보호는 거의 모든 카테고리 변수에 발생 — "
               "NaN 보존이 0 imputation 보다 항상 우선.",
               {"style": "subtitle", "color": C.on_dark, "line_spacing": 1.5})])

    add_pill_button(s, MARGIN_X + Inches(0.6), cy + ch - Inches(0.8),
                    Inches(3.7), Inches(0.5),
                    "Next  →  M3 · 시간대별 생활인구",
                    primary=True, color=C.canvas, text_color=ACCENT)


# ---------- main ----------
def main():
    prs = Presentation()
    setup_169(prs)

    builders = [
        slide_title,                 # 01
        slide_objectives,            # 02
        slide_data_meta,             # 03
        slide_na_protection,         # 04
        slide_register_vs_sample,    # 05
        slide_var_groups,            # 06
        slide_pivot,                 # 07
        slide_visualizations,        # 08
        slide_na_analysis,           # 09
        slide_admdong_groupby,       # 10
        slide_key_message,           # 11
    ]
    global TOTAL
    TOTAL = len(builders)

    for i, fn in enumerate(builders, start=1):
        fn(prs, i)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "..", "slides", "02_sgis_census.pptx")
    out = os.path.abspath(out)
    prs.save(out)
    print(f"OK · saved {out}  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
