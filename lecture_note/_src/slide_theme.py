"""MiniMax-inspired slide theme tokens & primitives for python-pptx.

16:9 only. White canvas with vibrant gradient cards, black pill CTAs.
Latin font: DM Sans (substitutes if absent). East Asian: Malgun Gothic.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------- canvas ----------
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.6)
MARGIN_Y = Inches(0.55)
CONTENT_W = SLIDE_W - MARGIN_X * 2
CONTENT_H = SLIDE_H - MARGIN_Y * 2


# ---------- colors (RGB) ----------
class C:
    # brand product colors — confined to product-identity moments
    coral        = RGBColor(0xFF, 0x4E, 0x3D)
    magenta      = RGBColor(0xE5, 0x1A, 0x7D)
    blue         = RGBColor(0x2B, 0x6C, 0xFF)
    blue_deep    = RGBColor(0x1D, 0x4E, 0xD8)
    blue_700     = RGBColor(0x1E, 0x40, 0xAF)
    cyan         = RGBColor(0x4D, 0xB6, 0xFF)
    blue_200     = RGBColor(0xDB, 0xEA, 0xFE)
    purple       = RGBColor(0x6D, 0x3A, 0xC9)

    # surfaces
    canvas       = RGBColor(0xFF, 0xFF, 0xFF)
    surface      = RGBColor(0xF6, 0xF6, 0xF7)
    surface_soft = RGBColor(0xFA, 0xFA, 0xFB)
    hairline     = RGBColor(0xE5, 0xE5, 0xE7)
    hairline_soft= RGBColor(0xEF, 0xEF, 0xF1)

    # text
    primary      = RGBColor(0x0A, 0x0A, 0x0A)   # near-black anchor
    ink_strong   = RGBColor(0x00, 0x00, 0x00)
    ink          = RGBColor(0x12, 0x12, 0x12)
    charcoal     = RGBColor(0x2C, 0x2C, 0x2E)
    slate        = RGBColor(0x57, 0x57, 0x60)
    steel        = RGBColor(0x80, 0x80, 0x8B)
    stone        = RGBColor(0xA1, 0xA1, 0xAA)
    muted        = RGBColor(0xB5, 0xB5, 0xBC)
    on_dark      = RGBColor(0xFF, 0xFF, 0xFF)
    on_primary   = RGBColor(0xFF, 0xFF, 0xFF)

    # semantic
    success_bg   = RGBColor(0xE7, 0xF7, 0xEC)
    success_text = RGBColor(0x18, 0x6A, 0x3E)
    warn_bg      = RGBColor(0xFD, 0xF1, 0xE3)
    warn_text    = RGBColor(0xB4, 0x53, 0x09)


# ---------- font ----------
LATIN_FONT = "DM Sans"
ASIAN_FONT = "Malgun Gothic"


def _set_font(run, *, size_pt, color, bold=False, latin=LATIN_FONT, asian=ASIAN_FONT,
              letter_spacing_pt: float | None = None, line_spacing: float | None = None):
    """Apply Latin + East Asian font + size + color in one shot."""
    run.font.size = Pt(size_pt)
    run.font.bold = bool(bold)
    run.font.color.rgb = color
    run.font.name = latin

    rPr = run._r.get_or_add_rPr()
    # East Asian font (Korean) — separate from Latin
    for tag in ("ea",):
        old = rPr.find(qn(f"a:{tag}"))
        if old is not None:
            rPr.remove(old)
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", asian)
    # ensure Latin too (overrides default theme face)
    latin_el = rPr.find(qn("a:latin"))
    if latin_el is None:
        latin_el = etree.SubElement(rPr, qn("a:latin"))
    latin_el.set("typeface", latin)

    # letter spacing (in 1/100 of a point)
    if letter_spacing_pt is not None:
        rPr.set("spc", str(int(letter_spacing_pt * 100)))


def _set_para(p, *, align=PP_ALIGN.LEFT, line_spacing: float | None = None,
              space_before: float | None = None, space_after: float | None = None):
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)


# ---------- typography presets (size, weight, line-height, letter-spacing) ----------
TYPE = {
    "hero":         dict(size_pt=60, bold=True,  line_spacing=1.10, letter_spacing_pt=-1.5),
    "display_lg":   dict(size_pt=42, bold=True,  line_spacing=1.10, letter_spacing_pt=-1.0),
    "heading_lg":   dict(size_pt=30, bold=True,  line_spacing=1.20, letter_spacing_pt=-0.5),
    "heading_md":   dict(size_pt=24, bold=True,  line_spacing=1.25, letter_spacing_pt=-0.3),
    "heading_sm":   dict(size_pt=18, bold=True,  line_spacing=1.30),
    "card_title":   dict(size_pt=15, bold=True,  line_spacing=1.35),
    "subtitle":     dict(size_pt=14, bold=False, line_spacing=1.50),
    "body_md":      dict(size_pt=12, bold=False, line_spacing=1.50),
    "body_md_b":    dict(size_pt=12, bold=True,  line_spacing=1.50),
    "body_sm":      dict(size_pt=11, bold=False, line_spacing=1.50),
    "body_sm_med":  dict(size_pt=11, bold=True,  line_spacing=1.50),
    "caption":      dict(size_pt=10, bold=False, line_spacing=1.50),
    "caption_b":    dict(size_pt=10, bold=True,  line_spacing=1.40),
    "micro":        dict(size_pt=9,  bold=False, line_spacing=1.40),
    "button":       dict(size_pt=11, bold=True,  line_spacing=1.30),
    "code":         dict(size_pt=11, bold=False, line_spacing=1.20),
}


# ---------- code block helper ----------
def add_code_block(slide, x, y, w, h, lines, *,
                   dark=True, label=None, show_dots=True,
                   font_pt=11, line_spacing=1.25,
                   padding_in=0.32):
    """Code panel.

    `lines` items can be:
      - str                       : single code line
      - (code, comment)           : code + trailing comment (two-tone)

    `dark=True`  → near-black panel, light text (IDE look)
    `dark=False` → surface panel, ink text (print look)
    """
    if dark:
        bg            = RGBColor(0x10, 0x14, 0x1F)
        code_color    = RGBColor(0xEC, 0xEE, 0xF2)
        comment_color = RGBColor(0x70, 0x77, 0x83)
        label_color   = RGBColor(0x97, 0x9C, 0xA6)
        divider_color = RGBColor(0x23, 0x28, 0x33)
    else:
        bg            = C.surface
        code_color    = C.ink
        comment_color = C.steel
        label_color   = C.steel
        divider_color = C.hairline_soft

    add_rounded(slide, x, y, w, h, fill=bg, line=None,
                radius_ratio=radius_for(w, h, 0.18))

    has_header = bool(label) or show_dots
    header_h = Inches(0.42) if has_header else Inches(0)

    if has_header:
        if show_dots:
            dot_d = Inches(0.13)
            dot_y = y + Inches(0.155)
            for i, dc in enumerate([RGBColor(0xFF, 0x5F, 0x57),
                                    RGBColor(0xFE, 0xBC, 0x2E),
                                    RGBColor(0x28, 0xC8, 0x40)]):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    x + Inches(padding_in + 0.18 * i),
                    dot_y, dot_d, dot_d,
                )
                dot.fill.solid(); dot.fill.fore_color.rgb = dc
                dot.line.fill.background()
                # remove default shadow
                spPr = dot._element.spPr
                for tag in ("a:effectLst", "a:effectDag"):
                    for el in spPr.findall(qn(tag)):
                        spPr.remove(el)
                etree.SubElement(spPr, qn("a:effectLst"))
        if label:
            add_text(slide, x + Inches(padding_in),
                     y + Inches(0.08), w - Inches(padding_in * 2),
                     Inches(0.3),
                     [(label, {"style": "caption_b", "color": label_color,
                               "letter_spacing_pt": 1.4})],
                     align=PP_ALIGN.RIGHT)
        # divider under header
        add_hairline(slide, x + Inches(padding_in), y + header_h,
                     w - Inches(padding_in * 2),
                     color=divider_color, weight_pt=0.5)

    body_x = x + Inches(padding_in)
    body_y = y + header_h + Inches(0.18 if has_header else padding_in)
    body_w = w - Inches(padding_in * 2)
    body_h = h - (y + h - body_y) * 0  # use rest

    tb = slide.shapes.add_textbox(body_x, body_y, body_w,
                                   y + h - body_y - Inches(padding_in))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(lines):
        if isinstance(item, str):
            code_str, comment_str = item, None
        else:
            code_str, comment_str = item

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para(p, align=PP_ALIGN.LEFT, line_spacing=line_spacing,
                  space_before=0, space_after=0)

        if code_str:
            run = p.add_run()
            run.text = code_str
            _set_font(run, size_pt=font_pt, color=code_color, bold=False,
                      latin="Consolas", asian=ASIAN_FONT)
        if comment_str:
            run2 = p.add_run()
            run2.text = ("  " if code_str else "") + comment_str
            _set_font(run2, size_pt=font_pt, color=comment_color, bold=False,
                      latin="Consolas", asian=ASIAN_FONT)
    return tb


# ---------- primitive: text in a frame ----------
def add_text(slide, x, y, w, h, lines, *,
             default_style="body_md", default_color=None, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, fill=None, padding=0):
    """Add a textbox. `lines` is list of (text, style_overrides) or plain str.

    style_overrides may include: style (preset key), size_pt, bold, color, latin,
    letter_spacing_pt, align, space_before, space_after, line_spacing, code (bool).
    """
    if default_color is None:
        default_color = C.ink

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(padding)
    tf.margin_right = Emu(padding)
    tf.margin_top = Emu(padding)
    tf.margin_bottom = Emu(padding)
    tf.vertical_anchor = anchor

    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
        tb.line.fill.background()
    else:
        tb.fill.background()
        tb.line.fill.background()

    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, ov = item, {}
        else:
            text, ov = item

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        style_key = ov.get("style", default_style)
        style = dict(TYPE[style_key])
        for k in ("size_pt", "bold", "letter_spacing_pt", "line_spacing"):
            if k in ov:
                style[k] = ov[k]
        color = ov.get("color", default_color)
        latin = ov.get("latin", LATIN_FONT)
        asian = ov.get("asian", ASIAN_FONT)
        if ov.get("code"):
            latin = "Consolas"

        _set_para(p,
                  align=ov.get("align", align),
                  line_spacing=style.get("line_spacing"),
                  space_before=ov.get("space_before"),
                  space_after=ov.get("space_after"))

        run = p.add_run()
        run.text = text
        _set_font(run,
                  size_pt=style["size_pt"],
                  color=color,
                  bold=style.get("bold", False),
                  latin=latin, asian=asian,
                  letter_spacing_pt=style.get("letter_spacing_pt"))
    return tb


# ---------- primitive: rounded rectangle ----------
def add_rounded(slide, x, y, w, h, *, fill=None, line=None, line_w_pt=1.0,
                radius_ratio=0.05, shadow=False):
    """Rounded rect. radius_ratio: corner softening as fraction of min dim."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # control adj1: 0.0–0.5 of half min dim
    shp.adjustments[0] = max(0.001, min(0.5, radius_ratio))

    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill

    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w_pt)

    # remove default shadow
    sp = shp.shadow
    sp.inherit = False
    if not shadow:
        # python-pptx exposes no shadow-off; suppress via XML
        spPr = shp._element.spPr
        for tag in ("a:effectLst", "a:effectDag"):
            for el in spPr.findall(qn(tag)):
                spPr.remove(el)
        effect = etree.SubElement(spPr, qn("a:effectLst"))

    shp.text_frame.text = ""
    return shp


# ---------- corner radius helpers ----------
def radius_for(w_emu, h_emu, target_in):
    """Convert target corner radius (inches) to adjustment ratio.

    Rounded rect adj1 = radius / (min_side / 2). 0.5 = full pill.
    """
    min_side = min(w_emu, h_emu)
    target_emu = Inches(target_in)
    return min(0.5, max(0.001, target_emu / (min_side / 2)))


def add_card_white(slide, x, y, w, h, *, radius_in=0.21):
    """Standard documentation card: white + hairline border + 16px (~0.21in) corners."""
    return add_rounded(slide, x, y, w, h,
                       fill=C.canvas, line=C.hairline, line_w_pt=0.75,
                       radius_ratio=radius_for(w, h, radius_in))


def add_card_surface(slide, x, y, w, h, *, radius_in=0.21):
    return add_rounded(slide, x, y, w, h,
                       fill=C.surface, line=None,
                       radius_ratio=radius_for(w, h, radius_in))


def add_card_vibrant(slide, x, y, w, h, *, color, radius_in=0.42):
    """Vibrant gradient/product card: 32px corners, no border."""
    return add_rounded(slide, x, y, w, h,
                       fill=color, line=None,
                       radius_ratio=radius_for(w, h, radius_in))


def add_pill_button(slide, x, y, w, h, label, *,
                    primary=True, color=None, text_color=None):
    """Black-pill primary or outline-pill secondary."""
    if color is None:
        color = C.primary if primary else None
    if text_color is None:
        text_color = C.on_primary if primary else C.ink

    if primary:
        shp = add_rounded(slide, x, y, w, h, fill=color, line=None,
                          radius_ratio=0.5)
    else:
        shp = add_rounded(slide, x, y, w, h, fill=None, line=C.ink,
                          line_w_pt=1.25, radius_ratio=0.5)

    tf = shp.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _set_para(p, align=PP_ALIGN.CENTER, line_spacing=1.0)
    run = p.add_run()
    run.text = label
    _set_font(run, size_pt=TYPE["button"]["size_pt"], color=text_color, bold=True)
    return shp


def add_badge(slide, x, y, label, *, kind="new"):
    """Pill badge. kind: 'new' (coral), 'beta' (blue-200), 'success' (green), 'code' (blue-200 sm)."""
    palette = {
        "new":     (C.coral,        C.on_dark,      0.5),
        "success": (C.success_bg,   C.success_text, 0.5),
        "beta":    (C.blue_200,     C.blue_deep,    0.5),
        "code":    (C.blue_200,     C.blue_deep,    0.18),
    }[kind]
    fill, text_color, radius_ratio = palette
    # auto width based on label length (approx)
    w = Inches(0.25 + 0.085 * max(len(label), 3))
    h = Inches(0.27)
    shp = add_rounded(slide, x, y, w, h, fill=fill, line=None,
                      radius_ratio=radius_ratio)
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _set_para(p, align=PP_ALIGN.CENTER, line_spacing=1.0)
    run = p.add_run()
    run.text = label
    _set_font(run, size_pt=9, color=text_color, bold=True, letter_spacing_pt=0.4)
    return shp


# ---------- primitive: hairline divider ----------
def add_hairline(slide, x, y, w, *, color=None, weight_pt=0.75):
    if color is None:
        color = C.hairline
    line = slide.shapes.add_connector(1, x, y, x + w, y)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


# ---------- chrome: page header & footer ----------
def add_page_chrome(slide, *, module_label, page_no, total_pages, deck_label="MiniMax Studio  ·  TAZ Lecture"):
    """Top: thin module/section label. Bottom: deck label + page x/y."""
    # top-left module label
    add_text(slide, MARGIN_X, Inches(0.28), Inches(8), Inches(0.3),
             [(module_label, {"style": "caption_b", "color": C.steel, "letter_spacing_pt": 0.6})])
    # top-right tiny brand chip
    add_text(slide, SLIDE_W - MARGIN_X - Inches(2.5), Inches(0.28),
             Inches(2.5), Inches(0.3),
             [("MiniMax Studio", {"style": "caption_b", "color": C.ink, "letter_spacing_pt": 0.4})],
             align=PP_ALIGN.RIGHT)
    # bottom deck label + page
    add_hairline(slide, MARGIN_X, SLIDE_H - Inches(0.55),
                 SLIDE_W - MARGIN_X * 2, color=C.hairline_soft)
    add_text(slide, MARGIN_X, SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             [(deck_label, {"style": "micro", "color": C.steel})])
    add_text(slide, SLIDE_W - MARGIN_X - Inches(2), SLIDE_H - Inches(0.45),
             Inches(2), Inches(0.3),
             [(f"{page_no:02d} / {total_pages:02d}",
               {"style": "micro", "color": C.steel, "letter_spacing_pt": 0.5})],
             align=PP_ALIGN.RIGHT)


# ---------- helpers ----------
def new_blank_slide(prs):
    """Add a blank-layout slide with white canvas."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C.canvas
    return slide


def setup_169(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
