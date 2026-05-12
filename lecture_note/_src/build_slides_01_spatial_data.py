"""01_spatial_data.pptx — M1 슬라이드 (16:9, MiniMax 디자인 시스템)."""
from __future__ import annotations

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from slide_theme import (
    C, SLIDE_W, SLIDE_H, MARGIN_X, MARGIN_Y, CONTENT_W, CONTENT_H,
    add_text, add_rounded, add_card_white, add_card_surface, add_card_vibrant,
    add_pill_button, add_badge, add_hairline, add_page_chrome,
    add_code_block,
    new_blank_slide, setup_169, radius_for, TYPE,
)

MODULE_LABEL = "M1  ·  서울 공간 데이터 둘러보기"
DECK_LABEL = "MiniMax Studio  ·  TAZ Lecture  ·  M1"
TOTAL = 13   # set after composition; updated at end if needed


# ---------- slide 1: hero title ----------
def slide_title(prs, idx):
    s = new_blank_slide(prs)

    # full-bleed coral accent on the right ~38%
    coral_w = Inches(5.0)
    add_rounded(s, SLIDE_W - coral_w - Inches(0.4), Inches(0.4),
                coral_w, SLIDE_H - Inches(0.8),
                fill=C.coral, line=None,
                radius_ratio=radius_for(coral_w, SLIDE_H - Inches(0.8), 0.42))

    # promo strip at very top (black)
    promo_w = Inches(4.5)
    promo = add_rounded(s, MARGIN_X, Inches(0.45), promo_w, Inches(0.42),
                        fill=C.primary, line=None, radius_ratio=0.5)
    tf = promo.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    from slide_theme import _set_para, _set_font
    _set_para(p, align=PP_ALIGN.CENTER, line_spacing=1.0)
    r = p.add_run(); r.text = "LECTURE  ·  M1 of 7  ·  Spatial Data"
    _set_font(r, size_pt=10, color=C.on_primary, bold=True, letter_spacing_pt=0.8)

    # hero
    add_text(s, MARGIN_X, Inches(1.7), Inches(7.2), Inches(2.6),
             [("서울 공간\n데이터 둘러보기",
               {"style": "hero", "color": C.ink, "letter_spacing_pt": -1.8, "line_spacing": 1.05})])

    # subtitle
    add_text(s, MARGIN_X, Inches(4.6), Inches(7.2), Inches(1.0),
             [("집계구 · 행정동 · 도로 · 필지의 위계와 좌표계.",
               {"style": "subtitle", "color": C.slate})])

    # buttons
    btn_y = Inches(5.55)
    add_pill_button(s, MARGIN_X, btn_y, Inches(2.0), Inches(0.5),
                    "Open Notebook  →", primary=True)
    add_pill_button(s, MARGIN_X + Inches(2.15), btn_y, Inches(1.7), Inches(0.5),
                    "View Repo", primary=False)

    # right card content (on coral)
    card_x = SLIDE_W - coral_w - Inches(0.0)  # text inside coral block
    inner_x = SLIDE_W - coral_w - Inches(0.0)
    pad = Inches(0.45)
    # M1 wordmark big
    add_text(s, SLIDE_W - coral_w - Inches(0.0) + pad, Inches(0.95),
             coral_w - pad * 2, Inches(2.4),
             [("M1", {"style": "hero", "color": C.on_dark,
                       "letter_spacing_pt": -2.5, "line_spacing": 1.0,
                       "size_pt": 96})])
    # tagline
    add_text(s, SLIDE_W - coral_w - Inches(0.0) + pad, Inches(3.5),
             coral_w - pad * 2, Inches(0.5),
             [("SPATIAL DATA", {"style": "caption_b",
                                "color": C.on_dark, "letter_spacing_pt": 1.5})])
    # bullet list
    bullets = [
        ("• 19,097 개 집계구 (OA)", {"style": "body_md", "color": C.on_dark}),
        ("• 414,791 개 도로 centerline", {"style": "body_md", "color": C.on_dark, "space_before": 4}),
        ("• 89.9 만 LSMD 필지", {"style": "body_md", "color": C.on_dark, "space_before": 4}),
        ("• 11 zone class 매핑", {"style": "body_md", "color": C.on_dark, "space_before": 4}),
    ]
    add_text(s, SLIDE_W - coral_w - Inches(0.0) + pad, Inches(4.2),
             coral_w - pad * 2, Inches(2.6), bullets)

    # NEW badge in card
    add_badge(s, SLIDE_W - coral_w + pad, Inches(0.6), "M1 · 2026", kind="new")

    # bottom chrome
    add_text(s, MARGIN_X, SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             [(DECK_LABEL, {"style": "micro", "color": C.steel})])
    add_text(s, SLIDE_W - MARGIN_X - Inches(2), SLIDE_H - Inches(0.45),
             Inches(2), Inches(0.3),
             [(f"{idx:02d} / {TOTAL:02d}",
               {"style": "micro", "color": C.steel, "letter_spacing_pt": 0.5})],
             align=PP_ALIGN.RIGHT)


# ---------- slide 2: 학습 목표 ----------
def slide_objectives(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(1.2),
             [("학습 목표",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.0), Inches(11), Inches(0.5),
             [("Five things you'll be able to do by the end of this module.",
               {"style": "subtitle", "color": C.slate})])

    objectives = [
        ("01", "SGIS 집계구(OA) 코드 체계와 위계 이해",
         "시군구 → 행정동 → 집계구 → 필지의 자릿수와 키 관계."),
        ("02", "EPSG 5179 / 5186 / 4326 좌표계 변환",
         "분석 표준은 5179. 면적·거리는 m 단위 CRS 에서만 의미."),
        ("03", "도로 두 SHP 의 역할 분담",
         "segment.ROA_CLS_SE = 등급 권위, centerline = link 그래프·물리속성."),
        ("04", "LSMD 필지 + 11 zone class 매핑",
         "ZONE_CLASS_MAP 단일 진실. 합 ≈ 605 km²."),
        ("05", "OA + 도로(2종) + 필지 3-layer overlay",
         "한 동 zoom-in 으로 같은 bbox 에서 비교."),
    ]

    # 5 cards in row with wrap (3 + 2)
    card_w = (CONTENT_W - Inches(0.4) * 2) / 3
    card_h = Inches(1.85)
    gap = Inches(0.25)
    start_y = Inches(2.85)

    for i, (no, title, desc) in enumerate(objectives):
        col = i % 3
        row = i // 3
        x = MARGIN_X + (card_w + gap) * col
        y = start_y + (card_h + gap) * row
        add_card_white(s, x, y, card_w, card_h, radius_in=0.21)
        # number tag
        add_text(s, x + Inches(0.25), y + Inches(0.2),
                 Inches(0.7), Inches(0.3),
                 [(no, {"style": "caption_b", "color": C.coral, "letter_spacing_pt": 0.8})])
        # title
        add_text(s, x + Inches(0.25), y + Inches(0.5),
                 card_w - Inches(0.5), Inches(0.55),
                 [(title, {"style": "card_title", "color": C.ink})])
        # desc
        add_text(s, x + Inches(0.25), y + Inches(1.05),
                 card_w - Inches(0.5), Inches(0.7),
                 [(desc, {"style": "body_sm", "color": C.slate})])


# ---------- slide 3: 공간 단위 위계 ----------
def slide_hierarchy(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(1.2),
             [("공간 단위 위계",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.0), Inches(11), Inches(0.5),
             [("같은 자릿수라도 체계가 다를 수 있다 — 행안부 vs SGIS.",
               {"style": "subtitle", "color": C.slate})])

    # left: code block (hierarchy tree)
    code_x = MARGIN_X
    code_y = Inches(2.85)
    code_w = Inches(7.0)
    code_h = Inches(3.7)

    tree_lines = [
        ("서울특별시", None),
        ("└─ 시군구 (25)",                  "5자리   · 예) 종로구 = 11010"),
        ("   └─ 행정동 (426)",              "8자리   · 시군구 + 3자리"),
        ("      └─ 집계구 (19,097)",        "14자리  · 행정동 + 일련번호 6"),
        ("         └─ 필지 (89.9 만)",      "19자리  · PNU"),
    ]
    add_code_block(s, code_x, code_y, code_w, code_h, tree_lines,
                   dark=True, label="seoul.hierarchy", show_dots=True,
                   font_pt=12, line_spacing=1.40, padding_in=0.36)

    # right: callout
    cx = code_x + code_w + Inches(0.3)
    cw = SLIDE_W - MARGIN_X - cx
    add_card_vibrant(s, cx, code_y, cw, code_h, color=C.coral, radius_in=0.42)
    add_text(s, cx + Inches(0.4), code_y + Inches(0.4),
             cw - Inches(0.8), Inches(0.4),
             [("핵심", {"style": "caption_b", "color": C.on_dark, "letter_spacing_pt": 1.2})])
    add_text(s, cx + Inches(0.4), code_y + Inches(0.85),
             cw - Inches(0.8), Inches(1.6),
             [("OA 는\n공식 최소 격자",
               {"style": "heading_md", "color": C.on_dark, "size_pt": 28,
                "line_spacing": 1.15, "letter_spacing_pt": -0.5})])
    add_text(s, cx + Inches(0.4), code_y + Inches(2.4),
             cw - Inches(0.8), Inches(1.2),
             [("통계청이 인구센서스 단위로 그어둔 약 500 세대 기준 폴리곤. "
               "본 프로젝트의 모든 통계는 이 키로 join 된다.",
               {"style": "body_sm", "color": C.on_dark, "line_spacing": 1.55})])


# ---------- slide 4: SGIS OA 메타 ----------
def slide_oa_meta(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    # eyebrow
    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("DATA · SGIS 집계구",
               {"style": "caption_b", "color": C.coral, "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("bnd_oa_11_2025_2Q.shp",
               {"style": "heading_lg", "color": C.ink,
                "code": True, "letter_spacing_pt": -0.5})])

    # left: meta table (key/value rows)
    rows = [
        ("출처",     "통계청 SGIS"),
        ("기준",     "2025-06-30 (2025년 2분기)"),
        ("CRS",      "EPSG:5179"),
        ("행 수",    "19,097 폴리곤"),
        ("면적",     "median 11,608 m² (≈ 100 × 100 m)"),
        ("총 면적",  "605.3 km² (서울시 전체)"),
        ("인코딩",   "UTF-8"),
    ]
    table_x = MARGIN_X
    table_y = Inches(2.55)
    table_w = Inches(7.5)
    row_h = Inches(0.46)

    add_card_white(s, table_x, table_y,
                   table_w, row_h * len(rows) + Inches(0.3),
                   radius_in=0.21)

    for i, (k, v) in enumerate(rows):
        ry = table_y + Inches(0.15) + row_h * i
        add_text(s, table_x + Inches(0.5), ry, Inches(2.0), row_h,
                 [(k, {"style": "body_sm_med", "color": C.steel})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, table_x + Inches(2.6), ry, table_w - Inches(2.9), row_h,
                 [(v, {"style": "body_md", "color": C.ink})],
                 anchor=MSO_ANCHOR.MIDDLE)
        if i < len(rows) - 1:
            add_hairline(s, table_x + Inches(0.5),
                         ry + row_h - Inches(0.01),
                         table_w - Inches(1.0),
                         color=C.hairline_soft, weight_pt=0.5)

    # right: warning callout
    cx = table_x + table_w + Inches(0.3)
    cw = SLIDE_W - MARGIN_X - cx
    cy = Inches(2.55)
    ch = Inches(2.7)
    add_card_white(s, cx, cy, cw, ch, radius_in=0.21)
    add_badge(s, cx + Inches(0.35), cy + Inches(0.3), "주의", kind="new")
    add_text(s, cx + Inches(0.35), cy + Inches(0.85),
             cw - Inches(0.7), Inches(1.0),
             [("8자리 코드,\n두 가지 체계",
               {"style": "heading_md", "color": C.ink,
                "line_spacing": 1.15, "letter_spacing_pt": -0.4})])
    add_text(s, cx + Inches(0.35), cy + Inches(2.0),
             cw - Inches(0.7), Inches(0.6),
             [("종로구 SGIS = 11010\n행안부 = 11110",
               {"style": "code", "color": C.charcoal, "code": True})])

    # columns table below
    cy2 = cy + ch + Inches(0.25)
    ch2 = Inches(1.3)
    add_card_surface(s, cx, cy2, cw, ch2, radius_in=0.16)
    add_text(s, cx + Inches(0.35), cy2 + Inches(0.18),
             cw - Inches(0.7), Inches(0.35),
             [("COLUMNS",
               {"style": "caption_b", "color": C.steel, "letter_spacing_pt": 1.0})])
    add_text(s, cx + Inches(0.35), cy2 + Inches(0.55),
             cw - Inches(0.7), Inches(0.7),
             [("BASE_DATE",  {"style": "code", "color": C.ink, "code": True}),
              ("ADM_CD       · 행정동 8자리",
               {"style": "code", "color": C.charcoal, "code": True, "space_before": 2}),
              ("TOT_OA_CD    · 집계구 14자리 = ADM_CD(8) + seq(6)",
               {"style": "code", "color": C.charcoal, "code": True, "space_before": 2})])


# ---------- slide 5: 좌표계 3종 ----------
def slide_crs(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(11), Inches(1.2),
             [("좌표계 (CRS) 3종",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.0), Inches(11), Inches(0.5),
             [("분석 표준은 EPSG:5179. m 단위 CRS 에서만 면적이 m² 가 된다.",
               {"style": "subtitle", "color": C.slate})])

    # 3 vibrant CRS tiles
    epsgs = [
        ("5179", "m",      "분석 표준",   "SGIS · 도로 · oa_master", C.blue,    True),
        ("5186", "m",      "LSMD raw",    "분석 시 5179 변환",        C.purple,  False),
        ("4326", "degree", "웹 뷰어",      "deck.gl 호환",              C.cyan,    False),
    ]
    tile_w = (CONTENT_W - Inches(0.5) * 2) / 3
    tile_h = Inches(2.6)
    ty = Inches(2.85)
    for i, (epsg, unit, role, use, color, badge) in enumerate(epsgs):
        x = MARGIN_X + (tile_w + Inches(0.5)) * i
        add_card_vibrant(s, x, ty, tile_w, tile_h, color=color, radius_in=0.42)
        # eyebrow
        add_text(s, x + Inches(0.4), ty + Inches(0.35),
                 tile_w - Inches(0.8), Inches(0.35),
                 [("EPSG", {"style": "caption_b", "color": C.on_dark,
                            "letter_spacing_pt": 1.5})])
        # huge number
        add_text(s, x + Inches(0.4), ty + Inches(0.7),
                 tile_w - Inches(0.8), Inches(1.2),
                 [(epsg, {"style": "hero", "color": C.on_dark,
                          "size_pt": 64, "letter_spacing_pt": -1.5,
                          "line_spacing": 1.0})])
        # unit chip
        add_text(s, x + Inches(0.4), ty + Inches(1.85),
                 tile_w - Inches(0.8), Inches(0.3),
                 [(f"단위 · {unit}", {"style": "caption_b", "color": C.on_dark,
                                       "letter_spacing_pt": 0.8})])
        # role + use
        add_text(s, x + Inches(0.4), ty + Inches(2.15),
                 tile_w - Inches(0.8), Inches(0.4),
                 [(f"{role} — {use}",
                   {"style": "body_sm", "color": C.on_dark})])

    # bottom code panel (dark IDE look)
    cy = ty + tile_h + Inches(0.3)
    ch = SLIDE_H - cy - Inches(0.7)
    py_lines = [
        ("gdf.crs",                        "# 현재 CRS"),
        ("gdf = gdf.to_crs(5179)",         "# 5179 (m) 로 변환"),
        ("gdf.geometry.area",              "# ⚠ 5179 에서만 m² 가 된다"),
    ]
    add_code_block(s, MARGIN_X, cy, CONTENT_W, ch, py_lines,
                   dark=True, label="PYTHON", show_dots=True,
                   font_pt=13, line_spacing=1.45, padding_in=0.4)


# ---------- slide 6: 도로 두 SHP 비교 ----------
def slide_roads_two(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("DATA · 도로 SHP",
               {"style": "caption_b", "color": C.coral, "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("두 SHP, 두 가지 역할",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.35), Inches(11), Inches(0.5),
             [("등급의 권위는 segment, 물리속성·그래프는 centerline.",
               {"style": "subtitle", "color": C.slate})])

    # two side-by-side comparison cards
    cw = (CONTENT_W - Inches(0.4)) / 2
    ch = Inches(3.85)
    cy = Inches(3.05)

    # LEFT — segment (magenta)
    lx = MARGIN_X
    add_card_vibrant(s, lx, cy, cw, ch, color=C.magenta, radius_in=0.42)
    add_text(s, lx + Inches(0.5), cy + Inches(0.45),
             cw - Inches(1.0), Inches(0.35),
             [("ROLE · 경계 권위",
               {"style": "caption_b", "color": C.on_dark, "letter_spacing_pt": 1.5})])
    add_text(s, lx + Inches(0.5), cy + Inches(0.85),
             cw - Inches(1.0), Inches(0.7),
             [("segment",
               {"style": "heading_lg", "color": C.on_dark,
                "code": True, "letter_spacing_pt": -0.5})])
    add_text(s, lx + Inches(0.5), cy + Inches(1.55),
             cw - Inches(1.0), Inches(0.4),
             [("TL_SPRD_MANAGE  ·  행안부 도로명주소",
               {"style": "body_sm", "color": C.on_dark})])

    # stats row
    stats_l = [
        ("66,698", "LineString"),
        ("ROA_CLS_SE", "등급 권위"),
        ("2026-01", "기준"),
    ]
    sw = (cw - Inches(1.0)) / 3
    for i, (n, lbl) in enumerate(stats_l):
        sx = lx + Inches(0.5) + sw * i
        add_text(s, sx, cy + Inches(2.15), sw, Inches(0.55),
                 [(n, {"style": "card_title", "color": C.on_dark,
                       "size_pt": 18, "code": (i==1), "letter_spacing_pt": -0.3})])
        add_text(s, sx, cy + Inches(2.7), sw, Inches(0.3),
                 [(lbl, {"style": "micro", "color": C.on_dark,
                          "letter_spacing_pt": 0.5})])

    # role tag
    add_text(s, lx + Inches(0.5), cy + ch - Inches(0.7),
             cw - Inches(1.0), Inches(0.5),
             [("→ build_blocks_oa.py · superblock 경계 polygonize",
               {"style": "body_sm_med", "color": C.on_dark})])

    # RIGHT — centerline (blue)
    rx = lx + cw + Inches(0.4)
    add_card_vibrant(s, rx, cy, cw, ch, color=C.blue, radius_in=0.42)
    add_text(s, rx + Inches(0.5), cy + Inches(0.45),
             cw - Inches(1.0), Inches(0.35),
             [("ROLE · 그래프 골격",
               {"style": "caption_b", "color": C.on_dark, "letter_spacing_pt": 1.5})])
    add_text(s, rx + Inches(0.5), cy + Inches(0.85),
             cw - Inches(1.0), Inches(0.7),
             [("centerline",
               {"style": "heading_lg", "color": C.on_dark,
                "code": True, "letter_spacing_pt": -0.5})])
    add_text(s, rx + Inches(0.5), cy + Inches(1.55),
             cw - Inches(1.0), Inches(0.4),
             [("N3L_A0020000  ·  국토지리정보원 NSDI",
               {"style": "body_sm", "color": C.on_dark})])

    stats_r = [
        ("414,791", "LineString (≈ 6×)"),
        ("RDLN/RVWD/ONSD", "물리속성"),
        ("2025-08", "기준"),
    ]
    for i, (n, lbl) in enumerate(stats_r):
        sx = rx + Inches(0.5) + sw * i
        add_text(s, sx, cy + Inches(2.15), sw, Inches(0.55),
                 [(n, {"style": "card_title", "color": C.on_dark,
                       "size_pt": 18, "code": (i==1), "letter_spacing_pt": -0.3})])
        add_text(s, sx, cy + Inches(2.7), sw, Inches(0.3),
                 [(lbl, {"style": "micro", "color": C.on_dark,
                          "letter_spacing_pt": 0.5})])

    add_text(s, rx + Inches(0.5), cy + ch - Inches(0.7),
             cw - Inches(1.0), Inches(0.5),
             [("→ link_to_block · star_network · viewer ROAD mode",
               {"style": "body_sm_med", "color": C.on_dark})])


# ---------- slide 7: segment ROA_CLS_SE ----------
def slide_segment_grades(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("ROAD · segment", {"style": "caption_b", "color": C.coral,
                                  "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("ROA_CLS_SE — 도로 등급",
               {"style": "display_lg", "color": C.ink, "code": True})])
    add_text(s, MARGIN_X, Inches(2.4), Inches(11), Inches(0.5),
             [("등급 1~3 을 polygonize 해서 superblock 경계를 만든다 — M5 에서 자세히.",
               {"style": "subtitle", "color": C.slate})])

    grades = [
        ("1", "고속",                  "15",     C.coral),
        ("2", "국도",                  "523",    C.magenta),
        ("3", "지방·특·광역시도",       "8,504",  C.blue),
        ("4", "시·군·구도",            "57,656", C.steel),  # not in superblock; muted
    ]
    tw = (CONTENT_W - Inches(0.4) * 3) / 4
    th = Inches(2.65)
    ty = Inches(3.15)
    for i, (lvl, name, count, color) in enumerate(grades):
        x = MARGIN_X + (tw + Inches(0.4)) * i
        is_super = (i < 3)
        if is_super:
            add_card_vibrant(s, x, ty, tw, th, color=color, radius_in=0.42)
            text_color = C.on_dark
            sub_color = C.on_dark
        else:
            add_card_white(s, x, ty, tw, th, radius_in=0.21)
            text_color = C.ink
            sub_color = C.slate

        # level number
        add_text(s, x + Inches(0.35), ty + Inches(0.35),
                 tw - Inches(0.7), Inches(0.4),
                 [("LEVEL", {"style": "caption_b", "color": text_color,
                              "letter_spacing_pt": 1.2})])
        add_text(s, x + Inches(0.35), ty + Inches(0.7),
                 tw - Inches(0.7), Inches(1.0),
                 [(lvl, {"style": "hero", "color": text_color,
                         "size_pt": 72, "letter_spacing_pt": -1.5,
                         "line_spacing": 1.0})])
        add_text(s, x + Inches(0.35), ty + Inches(1.65),
                 tw - Inches(0.7), Inches(0.45),
                 [(name, {"style": "card_title", "color": text_color})])
        add_text(s, x + Inches(0.35), ty + Inches(2.1),
                 tw - Inches(0.7), Inches(0.45),
                 [(f"{count} 행",
                   {"style": "body_sm", "color": sub_color})])


# ---------- slide 8: centerline 메타 ----------
def slide_centerline_meta(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("ROAD · centerline",
               {"style": "caption_b", "color": C.blue_deep, "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("N3L_A0020000_11.shp",
               {"style": "display_lg", "color": C.ink, "code": True,
                "letter_spacing_pt": -0.5})])

    # left: meta table
    rows = [
        ("출처",     "국토지리정보원 연속수치지형도 (NSDI)"),
        ("기준",     "2025-08"),
        ("CRS",      "EPSG:5179  (.prj deprecated · set_crs allow_override)"),
        ("행 수",    "414,791 LineString  (segment 의 약 6배, 골목 포함)"),
    ]
    table_x = MARGIN_X
    table_y = Inches(2.6)
    table_w = Inches(7.5)
    row_h = Inches(0.55)

    add_card_white(s, table_x, table_y,
                   table_w, row_h * len(rows) + Inches(0.3),
                   radius_in=0.21)
    for i, (k, v) in enumerate(rows):
        ry = table_y + Inches(0.15) + row_h * i
        add_text(s, table_x + Inches(0.5), ry, Inches(2.0), row_h,
                 [(k, {"style": "body_sm_med", "color": C.steel})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, table_x + Inches(2.6), ry, table_w - Inches(2.9), row_h,
                 [(v, {"style": "body_md", "color": C.ink})],
                 anchor=MSO_ANCHOR.MIDDLE)
        if i < len(rows) - 1:
            add_hairline(s, table_x + Inches(0.5),
                         ry + row_h - Inches(0.01),
                         table_w - Inches(1.0),
                         color=C.hairline_soft, weight_pt=0.5)

    # right: pipeline callout
    cx = table_x + table_w + Inches(0.3)
    cw = SLIDE_W - MARGIN_X - cx
    cy = table_y
    ch = row_h * len(rows) + Inches(0.3)
    add_card_vibrant(s, cx, cy, cw, ch, color=C.blue, radius_in=0.42)
    add_text(s, cx + Inches(0.4), cy + Inches(0.35),
             cw - Inches(0.8), Inches(0.35),
             [("PIPELINE", {"style": "caption_b", "color": C.on_dark,
                            "letter_spacing_pt": 1.2})])
    add_text(s, cx + Inches(0.4), cy + Inches(0.75),
             cw - Inches(0.8), Inches(1.5),
             [("414k 링크가\n전부 여기서.",
               {"style": "heading_md", "color": C.on_dark, "size_pt": 26,
                "line_spacing": 1.15, "letter_spacing_pt": -0.4})])
    add_text(s, cx + Inches(0.4), cy + Inches(2.05),
             cw - Inches(0.8), Inches(0.6),
             [("link_to_block → star_network → viewer ROAD mode",
               {"style": "code", "color": C.on_dark, "code": True})])

    # bottom: hint about cols
    hy = table_y + ch + Inches(0.25)
    hh = SLIDE_H - hy - Inches(0.7)
    add_card_surface(s, MARGIN_X, hy, CONTENT_W, hh, radius_in=0.16)
    add_text(s, MARGIN_X + Inches(0.4), hy + Inches(0.18),
             CONTENT_W - Inches(0.8), Inches(0.4),
             [("NEXT", {"style": "caption_b", "color": C.steel,
                        "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X + Inches(0.4), hy + Inches(0.55),
             CONTENT_W - Inches(0.8), Inches(0.55),
             [("16개 컬럼 중 분석에 쓰는 7개 — 다음 슬라이드에서 자세히.",
               {"style": "body_md", "color": C.charcoal})])


# ---------- slide 9: centerline 핵심 컬럼 ----------
def slide_centerline_cols(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(11), Inches(1.2),
             [("centerline — 핵심 컬럼 6개",
               {"style": "display_lg", "color": C.ink,
                "letter_spacing_pt": -0.5, "size_pt": 36})])
    add_text(s, MARGIN_X, Inches(1.85), Inches(11), Inches(0.5),
             [("16개 중 분석에 실제로 쓰는 6개. 등급은 viewer 색, 물리속성은 link 가중·시각화.",
               {"style": "subtitle", "color": C.slate})])

    cols = [
        ("SCLS",  "표준노드링크 등급",  "A0021 고속 / A0022 국도 / A00231 시도",
         "viewer 색 기준"),
        ("RDDV",  "도로 분류",          "RDD001 ~",
         "보조"),
        ("RDLN",  "차선 수",             "1 ~ 15",
         "link 가중·시각화"),
        ("RVWD",  "도로 폭 (m)",         "median 5.5 / max 120",
         "M5 polygonize 후 면적 보정"),
        ("ONSD",  "일방통행",            "ITH001 일방 ≈ 14%",
         "그래프 방향"),
        ("DVYN",  "중앙분리대",          "CSU001 유 ≈ 2%",
         "보조"),
    ]

    # 3 × 2 grid of card chips
    cw = (CONTENT_W - Inches(0.3) * 2) / 3
    ch = Inches(1.85)
    gy = Inches(2.75)
    for i, (name, meaning, values, use) in enumerate(cols):
        col = i % 3
        row = i // 3
        x = MARGIN_X + (cw + Inches(0.3)) * col
        y = gy + (ch + Inches(0.25)) * row
        add_card_white(s, x, y, cw, ch, radius_in=0.21)
        # name (code)
        add_text(s, x + Inches(0.3), y + Inches(0.25),
                 cw - Inches(0.6), Inches(0.5),
                 [(name, {"style": "card_title", "color": C.ink, "code": True,
                          "size_pt": 18, "letter_spacing_pt": -0.3})])
        add_badge(s, x + cw - Inches(1.05), y + Inches(0.27), "API", kind="code")
        # meaning
        add_text(s, x + Inches(0.3), y + Inches(0.78),
                 cw - Inches(0.6), Inches(0.35),
                 [(meaning, {"style": "body_sm_med", "color": C.charcoal})])
        # values
        add_text(s, x + Inches(0.3), y + Inches(1.1),
                 cw - Inches(0.6), Inches(0.35),
                 [(values, {"style": "code", "color": C.slate, "code": True,
                            "size_pt": 9})])
        # use
        add_text(s, x + Inches(0.3), y + Inches(1.4),
                 cw - Inches(0.6), Inches(0.35),
                 [(f"→ {use}",
                   {"style": "body_sm", "color": C.blue_deep})])


# ---------- slide 10: LSMD 11 zone class ----------
def slide_lsmd_zones(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("DATA · LSMD 필지",
               {"style": "caption_b", "color": C.purple, "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("11 zone class 매핑",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.4), Inches(11), Inches(0.5),
             [("89.9 만 필지 · 1.5 GB · ZONE_CLASS_MAP 이 단일 진실. 합 ≈ 605 km².",
               {"style": "subtitle", "color": C.slate})])

    rows = [
        ("전용주거",           "7.76",   "10,144"),
        ("일반주거_저밀(1종)", "100.01", "92,284"),
        ("일반주거_중밀(2종)", "146.74", "471,004"),
        ("일반주거_고밀(3종)", "89.31",  "141,597"),
        ("준주거",             "13.17",  "41,528"),
        ("중심·일반상업",      "21.31",  "55,282"),
        ("근린·유통상업",      "2.29",   "3,708"),
        ("일반·준공업",        "19.07",  "24,712"),
        ("보전·자연녹지",      "203.70", "57,943"),
        ("생산녹지",           "2.22",   "513"),
        ("전용공업",           "0",      "0"),
    ]

    # left: table card
    table_x = MARGIN_X
    table_y = Inches(3.05)
    table_w = Inches(7.8)
    row_h = Inches(0.32)
    table_h = row_h * (len(rows) + 1) + Inches(0.4)
    add_card_white(s, table_x, table_y, table_w, table_h, radius_in=0.21)

    # header
    hy = table_y + Inches(0.18)
    add_text(s, table_x + Inches(0.4), hy, Inches(3.5), row_h,
             [("CLASS", {"style": "caption_b", "color": C.steel,
                          "letter_spacing_pt": 1.0})],
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, table_x + Inches(4.3), hy, Inches(1.7), row_h,
             [("AREA (km²)", {"style": "caption_b", "color": C.steel,
                                "letter_spacing_pt": 1.0})],
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    add_text(s, table_x + Inches(6.05), hy, Inches(1.5), row_h,
             [("PARCELS", {"style": "caption_b", "color": C.steel,
                            "letter_spacing_pt": 1.0})],
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
    add_hairline(s, table_x + Inches(0.4),
                 hy + row_h + Inches(0.02),
                 table_w - Inches(0.8),
                 color=C.hairline, weight_pt=0.5)

    for i, (name, area, parcels) in enumerate(rows):
        ry = hy + row_h * (i + 1) + Inches(0.05)
        add_text(s, table_x + Inches(0.4), ry, Inches(3.5), row_h,
                 [(name, {"style": "body_sm", "color": C.ink})],
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, table_x + Inches(4.3), ry, Inches(1.7), row_h,
                 [(area, {"style": "body_sm", "color": C.charcoal})],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
        add_text(s, table_x + Inches(6.05), ry, Inches(1.5), row_h,
                 [(parcels, {"style": "body_sm", "color": C.charcoal})],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)

    # right: stat cards
    rx = table_x + table_w + Inches(0.3)
    rw = SLIDE_W - MARGIN_X - rx

    # purple total card
    sy = Inches(3.05)
    sh = Inches(2.3)
    add_card_vibrant(s, rx, sy, rw, sh, color=C.purple, radius_in=0.42)
    add_text(s, rx + Inches(0.4), sy + Inches(0.4),
             rw - Inches(0.8), Inches(0.35),
             [("TOTAL AREA", {"style": "caption_b", "color": C.on_dark,
                              "letter_spacing_pt": 1.2})])
    add_text(s, rx + Inches(0.4), sy + Inches(0.8),
             rw - Inches(0.8), Inches(1.0),
             [("≈ 605", {"style": "hero", "color": C.on_dark,
                          "size_pt": 64, "letter_spacing_pt": -1.5,
                          "line_spacing": 1.0})])
    add_text(s, rx + Inches(0.4), sy + Inches(1.85),
             rw - Inches(0.8), Inches(0.4),
             [("km² · SGIS 집계구 합과 일치",
               {"style": "body_sm", "color": C.on_dark})])

    # source path card
    py = sy + sh + Inches(0.25)
    ph = Inches(1.6)
    add_card_surface(s, rx, py, rw, ph, radius_in=0.16)
    add_text(s, rx + Inches(0.4), py + Inches(0.25),
             rw - Inches(0.8), Inches(0.35),
             [("SOURCE", {"style": "caption_b", "color": C.steel,
                          "letter_spacing_pt": 1.2})])
    add_text(s, rx + Inches(0.4), py + Inches(0.65),
             rw - Inches(0.8), Inches(0.85),
             [("data/raw/lsmd_zoning/",
               {"style": "code", "color": C.ink, "code": True, "size_pt": 11,
                "line_spacing": 1.20}),
              ("AL_D154_11_20260412/",
               {"style": "code", "color": C.charcoal, "code": True, "size_pt": 11,
                "line_spacing": 1.20})])


# ---------- slide 11: 노트북 흐름 ----------
def slide_notebook_flow(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    add_text(s, MARGIN_X, Inches(0.85), Inches(10), Inches(0.4),
             [("PRACTICE · 01_spatial_data.ipynb",
               {"style": "caption_b", "color": C.coral, "letter_spacing_pt": 1.2})])
    add_text(s, MARGIN_X, Inches(1.25), Inches(11), Inches(1.0),
             [("노트북 흐름 7단계",
               {"style": "display_lg", "color": C.ink})])
    add_text(s, MARGIN_X, Inches(2.4), Inches(11), Inches(0.5),
             [("로드 → 검증 → 시각화 → 종합 비교.",
               {"style": "subtitle", "color": C.slate})])

    steps = [
        ("01", "로드",      "SGIS OA · 도로 segment · centerline · 필지 sample"),
        ("02", "검증",      "CRS · 면적 · 행 수 · 시군구별 OA 개수"),
        ("03", "시각화",    "종로구 zoom-in (OA + 행정동 경계)"),
        ("04", "segment",   "등급별 색상 (간선 vs 시군구도)"),
        ("05", "centerline","SCLS 등급 + RDLN 차선 분포"),
        ("06", "필지",      "zone_class 색상 + 한 동 zoom-in"),
        ("07", "종합",      "segment vs centerline 2-패널 비교"),
    ]

    # 7 steps as a vertical-rhythm row of 7 cards
    sw = (CONTENT_W - Inches(0.2) * 6) / 7
    sh = Inches(3.0)
    sy = Inches(3.05)
    for i, (no, title, desc) in enumerate(steps):
        x = MARGIN_X + (sw + Inches(0.2)) * i
        # alternate: emphasize step 04 & 07 vibrant
        if no in ("04", "07"):
            add_card_vibrant(s, x, sy, sw, sh, color=C.coral if no == "04" else C.blue,
                             radius_in=0.32)
            tc = C.on_dark
            sub = C.on_dark
        else:
            add_card_white(s, x, sy, sw, sh, radius_in=0.18)
            tc = C.ink
            sub = C.slate

        add_text(s, x + Inches(0.2), sy + Inches(0.25),
                 sw - Inches(0.4), Inches(0.35),
                 [(no, {"style": "caption_b", "color": tc,
                         "letter_spacing_pt": 1.0})])
        add_text(s, x + Inches(0.2), sy + Inches(0.7),
                 sw - Inches(0.4), Inches(0.55),
                 [(title, {"style": "card_title", "color": tc,
                            "size_pt": 14})])
        add_text(s, x + Inches(0.2), sy + Inches(1.45),
                 sw - Inches(0.4), sh - Inches(1.6),
                 [(desc, {"style": "body_sm", "color": sub,
                           "size_pt": 10, "line_spacing": 1.4})])


# ---------- slide 12: 핵심 메시지 ----------
def slide_key_message(prs, idx):
    s = new_blank_slide(prs)
    add_page_chrome(s, module_label=MODULE_LABEL, page_no=idx, total_pages=TOTAL,
                    deck_label=DECK_LABEL)

    # full-width vibrant card
    cy = Inches(1.05)
    ch = Inches(5.6)
    add_card_vibrant(s, MARGIN_X, cy, CONTENT_W, ch, color=C.coral, radius_in=0.42)

    add_text(s, MARGIN_X + Inches(0.6), cy + Inches(0.55),
             CONTENT_W - Inches(1.2), Inches(0.4),
             [("KEY MESSAGE", {"style": "caption_b", "color": C.on_dark,
                                "letter_spacing_pt": 1.5})])

    add_text(s, MARGIN_X + Inches(0.6), cy + Inches(1.05),
             CONTENT_W - Inches(1.2), Inches(3.0),
             [("데이터를 합치기 전에\n각 데이터의 단위 · 키 · 좌표계를\n따로따로 확실히 알아두는 것 — 이것이 가장 중요하다.",
               {"style": "hero", "color": C.on_dark, "size_pt": 40,
                "letter_spacing_pt": -1.0, "line_spacing": 1.15})])

    add_text(s, MARGIN_X + Inches(0.6), cy + ch - Inches(1.4),
             CONTENT_W - Inches(1.2), Inches(0.5),
             [("자릿수 같은 코드도 체계가 다를 수 있고 (행안부 vs SGIS), "
               "면적·거리는 m 단위 CRS (5179) 에서만 의미가 있다.",
               {"style": "subtitle", "color": C.on_dark, "line_spacing": 1.5})])

    # next module pill
    add_pill_button(s, MARGIN_X + Inches(0.6), cy + ch - Inches(0.8),
                    Inches(3.5), Inches(0.5),
                    "Next  →  M2 · SGIS 인구·가구·주택",
                    primary=True, color=C.canvas, text_color=C.coral)


# ---------- main ----------
def main():
    prs = Presentation()
    setup_169(prs)

    builders = [
        slide_title,           # 01
        slide_objectives,      # 02
        slide_hierarchy,       # 03
        slide_oa_meta,         # 04
        slide_crs,             # 05
        slide_roads_two,       # 06
        slide_segment_grades,  # 07
        slide_centerline_meta, # 08
        slide_centerline_cols, # 09
        slide_lsmd_zones,      # 10
        slide_notebook_flow,   # 11
        slide_key_message,     # 12
    ]
    global TOTAL
    TOTAL = len(builders)

    for i, fn in enumerate(builders, start=1):
        fn(prs, i)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "..", "slides", "01_spatial_data.pptx")
    out = os.path.abspath(out)
    prs.save(out)
    print(f"OK · saved {out}  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
