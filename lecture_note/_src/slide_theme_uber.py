"""Uber-inspired slide theme tokens & primitives for python-pptx.

16:9 only. Black-and-white duet with grayscale; pill is the signature
interactive shape (999px). Pretendard family for all text. Editorial
illustrations and accent palette intentionally absent.

Pretendard weight families installed on Windows (verified):
  - Pretendard            (Regular / Bold pair)
  - Pretendard ExtraBold
  - Pretendard Medium
  - Pretendard ExtraLight
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree


# ---------- canvas (16:9 only) ----------
SLIDE_W   = Inches(13.333)
SLIDE_H   = Inches(7.5)
MARGIN_X  = Inches(0.6)
MARGIN_Y  = Inches(0.55)
CONTENT_W = SLIDE_W - MARGIN_X * 2
CONTENT_H = SLIDE_H - MARGIN_Y * 2


# ---------- colors ----------
class C:
    # brand / conversion
    primary       = RGBColor(0x00, 0x00, 0x00)   # Ink Black — every CTA, footer, dark band
    on_primary    = RGBColor(0xFF, 0xFF, 0xFF)
    black_elev    = RGBColor(0x28, 0x28, 0x28)   # hover on translucent pills
    surface_press = RGBColor(0xE2, 0xE2, 0xE2)   # pressed state for white pills

    # surfaces
    canvas        = RGBColor(0xFF, 0xFF, 0xFF)
    canvas_soft   = RGBColor(0xEF, 0xEF, 0xEF)   # chips, form rows, subtle pills
    canvas_softer = RGBColor(0xF3, 0xF3, 0xF3)   # nested inputs on white

    # text on light
    ink           = RGBColor(0x00, 0x00, 0x00)
    body          = RGBColor(0x5E, 0x5E, 0x5E)
    hairline_mid  = RGBColor(0x4B, 0x4B, 0x4B)
    mute          = RGBColor(0xAF, 0xAF, 0xAF)

    # text on dark
    on_dark       = RGBColor(0xFF, 0xFF, 0xFF)

    # link (browser blue — legal fine print only)
    link          = RGBColor(0x00, 0x00, 0xEE)

    # hairlines
    hairline      = RGBColor(0xE5, 0xE5, 0xE5)
    hairline_soft = RGBColor(0xEE, 0xEE, 0xEE)


# ---------- fonts (Pretendard) ----------
F_EXTRABOLD  = "Pretendard ExtraBold"
F_BOLD       = "Pretendard"           # use with bold=True for Bold weight
F_MEDIUM     = "Pretendard Medium"
F_EXTRALIGHT = "Pretendard ExtraLight"
F_MONO       = "Cascadia Code"        # Pretendard is not monospace; fallback Consolas
F_MONO_FALL  = "Consolas"


# ---------- typography presets ----------
# (size_pt, font_family, bold, line_spacing, letter_spacing_pt)
# Sentence-case is the voice; no all-caps display.
TYPE = {
    "display_xxl":  dict(size_pt=44, font=F_EXTRABOLD, bold=False, line_spacing=1.10, letter_spacing_pt=-1.0),
    "display_xl":   dict(size_pt=32, font=F_EXTRABOLD, bold=False, line_spacing=1.15, letter_spacing_pt=-0.6),
    "display_lg":   dict(size_pt=26, font=F_EXTRABOLD, bold=False, line_spacing=1.20, letter_spacing_pt=-0.4),
    "display_md":   dict(size_pt=20, font=F_BOLD,      bold=True,  line_spacing=1.25, letter_spacing_pt=-0.2),
    "display_sm":   dict(size_pt=16, font=F_BOLD,      bold=True,  line_spacing=1.30, letter_spacing_pt=0.0),
    "lead":         dict(size_pt=15, font=F_MEDIUM,    bold=False, line_spacing=1.50),
    "body_md":      dict(size_pt=13, font=F_MEDIUM,    bold=False, line_spacing=1.55),
    "body_md_b":    dict(size_pt=13, font=F_BOLD,      bold=True,  line_spacing=1.50),
    "body_sm":      dict(size_pt=11, font=F_MEDIUM,    bold=False, line_spacing=1.55),
    "body_sm_b":    dict(size_pt=11, font=F_BOLD,      bold=True,  line_spacing=1.45),
    "caption":      dict(size_pt=10, font=F_EXTRALIGHT,bold=False, line_spacing=1.45),
    "caption_b":    dict(size_pt=10, font=F_BOLD,      bold=True,  line_spacing=1.40),
    "micro":        dict(size_pt=9,  font=F_EXTRALIGHT,bold=False, line_spacing=1.35),
    "eyebrow":      dict(size_pt=11, font=F_BOLD,      bold=True,  line_spacing=1.20, letter_spacing_pt=2.4),
    "button":       dict(size_pt=12, font=F_MEDIUM,    bold=False, line_spacing=1.20),
    "button_lg":    dict(size_pt=14, font=F_MEDIUM,    bold=False, line_spacing=1.20),
    "code":         dict(size_pt=11, font=F_MONO,      bold=False, line_spacing=1.18),
    "code_sm":      dict(size_pt=10, font=F_MONO,      bold=False, line_spacing=1.15),
    "code_xs":      dict(size_pt=9,  font=F_MONO,      bold=False, line_spacing=1.15),
}


def _set_font(run, *, size_pt, color, bold=False,
              latin=F_BOLD, asian=None, letter_spacing_pt=None):
    """Apply Latin + East Asian font + size + color. Korean falls back to same family."""
    if asian is None:
        asian = latin
    run.font.size = Pt(size_pt)
    run.font.bold = bool(bold)
    run.font.color.rgb = color
    run.font.name = latin

    rPr = run._r.get_or_add_rPr()
    # East Asian face
    for tag in ("ea",):
        old = rPr.find(qn(f"a:{tag}"))
        if old is not None:
            rPr.remove(old)
        el = etree.SubElement(rPr, qn(f"a:{tag}"))
        el.set("typeface", asian)
    # Latin face — overrides theme default
    latin_el = rPr.find(qn("a:latin"))
    if latin_el is None:
        latin_el = etree.SubElement(rPr, qn("a:latin"))
    latin_el.set("typeface", latin)

    if letter_spacing_pt is not None:
        rPr.set("spc", str(int(letter_spacing_pt * 100)))


def _set_para(p, *, align=PP_ALIGN.LEFT, line_spacing=None,
              space_before=None, space_after=None):
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    if space_before is not None:
        p.space_before = Pt(space_before)
    if space_after is not None:
        p.space_after = Pt(space_after)


# ---------- text in a frame ----------
def add_text(slide, x, y, w, h, lines, *,
             default_style="body_md", default_color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             fill=None, padding=0):
    """Add a textbox.

    `lines` items are str or (text, overrides_dict). overrides keys:
      style, size_pt, bold, color, font, code (bool, → monospace),
      align, line_spacing, letter_spacing_pt, space_before, space_after.
    """
    if default_color is None:
        default_color = C.ink

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(padding); tf.margin_right = Emu(padding)
    tf.margin_top = Emu(padding);  tf.margin_bottom = Emu(padding)
    tf.vertical_anchor = anchor

    if fill is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = fill
        tb.line.fill.background()
    else:
        tb.fill.background()
        tb.line.fill.background()

    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, ov = item, {}
        else:
            text, ov = item

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        style_key = ov.get("style", default_style)
        style = dict(TYPE[style_key])
        for k in ("size_pt", "bold", "letter_spacing_pt", "line_spacing"):
            if k in ov:
                style[k] = ov[k]
        color = ov.get("color", default_color)
        latin = ov.get("font", style["font"])
        if ov.get("code"):
            latin = F_MONO

        _set_para(p,
                  align=ov.get("align", align),
                  line_spacing=style.get("line_spacing"),
                  space_before=ov.get("space_before"),
                  space_after=ov.get("space_after"))

        run = p.add_run()
        run.text = text
        _set_font(run,
                  size_pt=style["size_pt"],
                  color=color,
                  bold=style.get("bold", False),
                  latin=latin,
                  letter_spacing_pt=style.get("letter_spacing_pt"))
    return tb


# ---------- rounded rectangle ----------
def radius_for(w_emu, h_emu, target_in):
    """Convert a target corner radius (inches) → adj1 ratio (0..0.5)."""
    min_side = min(w_emu, h_emu)
    target_emu = Inches(target_in)
    return min(0.5, max(0.001, target_emu / (min_side / 2)))


def add_rounded(slide, x, y, w, h, *, fill=None, line=None,
                line_w_pt=0.75, radius_ratio=0.05, shadow=False):
    """Rounded rectangle. radius_ratio=0.5 == full pill (999 px)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = max(0.001, min(0.5, radius_ratio))

    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill

    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w_pt)

    # strip default shadow
    sp = shp.shadow; sp.inherit = False
    if not shadow:
        spPr = shp._element.spPr
        for tag in ("a:effectLst", "a:effectDag"):
            for el in spPr.findall(qn(tag)):
                spPr.remove(el)
        etree.SubElement(spPr, qn("a:effectLst"))

    shp.text_frame.text = ""
    return shp


# ---------- card primitives ----------
def add_card_content(slide, x, y, w, h, *, radius_in=0.21):
    """White card with hairline border — canonical content card (16 px corners)."""
    return add_rounded(slide, x, y, w, h,
                       fill=C.canvas, line=C.hairline, line_w_pt=0.75,
                       radius_ratio=radius_for(w, h, radius_in))


def add_card_soft(slide, x, y, w, h, *, radius_in=0.21):
    """Gray-tinted sub-region card (canvas-soft fill, no border)."""
    return add_rounded(slide, x, y, w, h,
                       fill=C.canvas_soft, line=None,
                       radius_ratio=radius_for(w, h, radius_in))


def add_card_dark(slide, x, y, w, h, *, radius_in=0.21):
    """Black polarity-flip promo card. White text inside."""
    return add_rounded(slide, x, y, w, h,
                       fill=C.primary, line=None,
                       radius_ratio=radius_for(w, h, radius_in))


# ---------- pill button ----------
def add_pill(slide, x, y, w, h, label, *,
             variant="primary", style="button"):
    """Pill (radius=999 px).

    variants:
      'primary'  — black fill, white text (conversion CTA)
      'secondary'— white fill, black text, hairline border (paired with primary)
      'subtle'   — canvas-soft fill, ink text (tertiary / chip-like)
      'on-dark'  — black fill, white text on dark band (download-app pill)
    """
    if variant == "primary":
        shp = add_rounded(slide, x, y, w, h, fill=C.primary, line=None,
                          radius_ratio=0.5)
        text_color = C.on_primary
    elif variant == "secondary":
        shp = add_rounded(slide, x, y, w, h, fill=C.canvas, line=C.ink,
                          line_w_pt=1.25, radius_ratio=0.5)
        text_color = C.ink
    elif variant == "subtle":
        shp = add_rounded(slide, x, y, w, h, fill=C.canvas_soft, line=None,
                          radius_ratio=0.5)
        text_color = C.ink
    elif variant == "on-dark":
        shp = add_rounded(slide, x, y, w, h, fill=C.primary, line=None,
                          radius_ratio=0.5)
        text_color = C.on_primary
    else:
        raise ValueError(f"unknown pill variant: {variant}")

    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _set_para(p, align=PP_ALIGN.CENTER, line_spacing=1.0)
    style_dict = TYPE[style]
    r = p.add_run(); r.text = label
    _set_font(r,
              size_pt=style_dict["size_pt"],
              color=text_color,
              bold=style_dict.get("bold", False),
              latin=style_dict["font"])
    return shp


def add_chip(slide, x, y, label, *, h_in=0.34, pad_x_in=0.55, style="body_sm_b"):
    """Auto-width subtle pill (category chip)."""
    # heuristic width: per-char ~ 0.085" + 2*pad
    w_in = pad_x_in * 2 + 0.085 * max(len(label), 3)
    w, h = Inches(w_in), Inches(h_in)
    shp = add_rounded(slide, x, y, w, h, fill=C.canvas_soft, line=None,
                      radius_ratio=0.5)
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _set_para(p, align=PP_ALIGN.CENTER, line_spacing=1.0)
    sd = TYPE[style]
    r = p.add_run(); r.text = label
    _set_font(r, size_pt=sd["size_pt"], color=C.ink,
              bold=sd.get("bold", False), latin=sd["font"])
    return shp, w


# ---------- code cell (light, tight, code-cell look) ----------
def add_code_cell(slide, x, y, w, h, lines, *,
                  font_pt=11, line_spacing=1.18,
                  label=None, padding_in=0.30,
                  comment_color=None, bg=None, code_color=None):
    """Light code cell on canvas-softer surface (16 px corners, no chrome dots).

    `lines` items can be:
      - str                   : single line
      - (code, comment)       : code + trailing comment (two-tone)
    """
    if bg is None:           bg = C.canvas_softer
    if code_color is None:   code_color = C.ink
    if comment_color is None:comment_color = C.body

    # cell
    add_rounded(slide, x, y, w, h, fill=bg, line=C.hairline_soft,
                line_w_pt=0.5,
                radius_ratio=radius_for(w, h, 0.18))

    head_h = Inches(0.34) if label else Inches(0.0)

    if label:
        add_text(slide,
                 x + Inches(padding_in), y + Inches(0.10),
                 w - Inches(padding_in * 2), Inches(0.28),
                 [(label, {"style": "caption_b", "color": C.body,
                           "letter_spacing_pt": 1.2})])
        # hairline divider below label
        add_hairline(slide,
                     x + Inches(padding_in),
                     y + head_h,
                     w - Inches(padding_in * 2),
                     color=C.hairline_soft, weight_pt=0.4)

    body_x = x + Inches(padding_in)
    body_y = y + head_h + Inches(0.10 if label else padding_in)
    body_w = w - Inches(padding_in * 2)

    tb = slide.shapes.add_textbox(body_x, body_y, body_w,
                                  y + h - body_y - Inches(padding_in))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(lines):
        if isinstance(item, str):
            code_str, comment_str = item, None
        else:
            code_str, comment_str = item

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para(p, align=PP_ALIGN.LEFT, line_spacing=line_spacing,
                  space_before=0, space_after=0)

        if code_str:
            run = p.add_run()
            run.text = code_str
            _set_font(run, size_pt=font_pt, color=code_color, bold=False,
                      latin=F_MONO, asian=F_MEDIUM)
        if comment_str:
            run2 = p.add_run()
            run2.text = ("  " if code_str else "") + comment_str
            _set_font(run2, size_pt=font_pt, color=comment_color, bold=False,
                      latin=F_MONO, asian=F_MEDIUM)
    return tb


# ---------- hairline divider ----------
def add_hairline(slide, x, y, w, *, color=None, weight_pt=0.6):
    if color is None:
        color = C.hairline
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x + w, y)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


def add_vline(slide, x, y, h, *, color=None, weight_pt=0.6):
    if color is None:
        color = C.hairline
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y, x, y + h)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line


# ---------- arrow connector ----------
def add_arrow(slide, x1, y1, x2, y2, *, color=None, weight_pt=1.25):
    if color is None:
        color = C.hairline_mid
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    # add arrowhead at endpoint
    lnPr = line.line._get_or_add_ln()
    # tailEnd / headEnd
    for tag in ("a:headEnd", "a:tailEnd"):
        for el in lnPr.findall(qn(tag)):
            lnPr.remove(el)
    head = etree.SubElement(lnPr, qn("a:tailEnd"))
    head.set("type", "triangle")
    head.set("w", "med"); head.set("len", "med")
    return line


# ---------- page chrome ----------
def add_page_chrome(slide, *, eyebrow, page_no, total_pages,
                    deck="TAZ · Lecture 1 · 서울 공간 데이터"):
    """Top: eyebrow (uppercase). Bottom: hairline + deck + page x/y."""
    add_text(slide, MARGIN_X, Inches(0.32), Inches(8), Inches(0.3),
             [(eyebrow, {"style": "eyebrow", "color": C.ink})])
    # top-right small mark
    add_text(slide, SLIDE_W - MARGIN_X - Inches(3), Inches(0.32),
             Inches(3), Inches(0.3),
             [("TAZ STUDIO", {"style": "eyebrow", "color": C.body,
                              "letter_spacing_pt": 2.4})],
             align=PP_ALIGN.RIGHT)
    # bottom hairline + meta
    add_hairline(slide, MARGIN_X, SLIDE_H - Inches(0.52),
                 SLIDE_W - MARGIN_X * 2, color=C.hairline_soft)
    add_text(slide, MARGIN_X, SLIDE_H - Inches(0.42), Inches(8), Inches(0.3),
             [(deck, {"style": "micro", "color": C.body})])
    add_text(slide, SLIDE_W - MARGIN_X - Inches(2), SLIDE_H - Inches(0.42),
             Inches(2), Inches(0.3),
             [(f"{page_no:02d} / {total_pages:02d}",
               {"style": "micro", "color": C.body, "letter_spacing_pt": 0.5})],
             align=PP_ALIGN.RIGHT)


# ---------- slide & deck helpers ----------
def new_blank_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C.canvas
    return slide


def new_dark_slide(prs):
    """Polarity-flipped slide — black canvas, white text."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = C.primary
    return slide


def setup_169(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
